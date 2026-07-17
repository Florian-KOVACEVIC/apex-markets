"""
Apex Markets — Application d'Analyse Financière
-----------------------------------------------
Auteur  : FloKov
Stack   : Streamlit · yfinance · Pandas · NumPy · Plotly
Usage   : streamlit run app2.py    
""" 

import json
import warnings
warnings.filterwarnings("ignore")

import streamlit as st
import yfinance as yf
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime, timedelta
import io
import urllib.parse
import urllib.request
from scipy import stats as scipy_stats

# ══════════════════════════════════════════════════════════════
#  TOOLTIPS — Définitions des métriques clés
# ══════════════════════════════════════════════════════════════

TOOLTIPS = {
    # Prix
    "Prix actuel":    "Dernier prix de transaction de l'actif sur le marché.",
    "Open":           "Prix d'ouverture de la séance en cours.",
    "High (jour)":    "Prix le plus haut atteint aujourd'hui.",
    "Low (jour)":     "Prix le plus bas atteint aujourd'hui.",
    "Clôture préc.":  "Prix de clôture de la séance précédente.",
    "52-sem. High":   "Plus haut prix atteint sur les 52 dernières semaines.",
    "52-sem. Low":    "Plus bas prix atteint sur les 52 dernières semaines.",
    "Volume":         "Nombre de titres échangés sur la séance, un volume élevé confirme la tendance.",
    # Valorisation
    "Market Cap":     "Capitalisation boursière = Prix × Nombre d'actions. Mesure la taille de l'entreprise.",
    "Enterprise V.":  "Valeur d'entreprise = Market Cap + Dette - Trésorerie. Prix réel d'une acquisition.",
    "Beta":           "Sensibilité au marché. Beta > 1 : plus volatile que le marché. Beta < 1 : plus stable.",
    "Float":          "Nombre d'actions librement négociables en bourse (hors actionnaires bloqués).",
    "P/E (TTM)":      "Price/Earnings sur 12 mois glissants. Combien paie-t-on pour 1€ de bénéfice ? Élevé = chère ou forte croissance attendue.",
    "P/E Forward":    "P/E basé sur les bénéfices futurs estimés par les analystes. Reflect les anticipations.",
    "PEG":            "P/E divisé par le taux de croissance. PEG < 1 = potentiellement sous-évalué. PEG > 2 = cher.",
    "P/B":            "Price/Book. Prix vs valeur comptable des actifs. P/B < 1 = se négocie sous sa valeur nette.",
    "P/S":            "Price/Sales. Valorisation rapportée au chiffre d'affaires. Utile quand il n'y a pas de bénéfices.",
    "EV/EBITDA":      "Valeur d'entreprise / EBITDA. Multiple de valorisation neutre à la structure financière. < 10 = raisonnable.",
    "EV/Rev.":        "Valeur d'entreprise / Chiffre d'affaires. Valorisation par rapport aux ventes.",
    "52w perf":       "Performance du prix sur 52 semaines glissantes.",
    # Rentabilité
    "ROE":            "Return On Equity. Bénéfice net / Capitaux propres. Mesure l'efficacité pour créer de la valeur. > 15% = bon.",
    "ROA":            "Return On Assets. Bénéfice net / Total actifs. Efficacité d'utilisation de tous les actifs.",
    "Marge brute":    "Chiffre d'affaires - Coût des ventes. Mesure la rentabilité avant les frais généraux.",
    "Marge opérat.":  "Bénéfice opérationnel / CA. Rentabilité après coûts d'exploitation, avant impôts et intérêts.",
    "Marge nette":    "Bénéfice net / CA. Ce qui reste après toutes les charges. Mesure la rentabilité finale.",
    "EBITDA":         "Bénéfices avant intérêts, impôts, amortissements. Proxy du cash-flow opérationnel.",
    "Chiffre d'aff.": "Revenus totaux générés par l'entreprise sur 12 mois.",
    "Rev. Growth":    "Taux de croissance du chiffre d'affaires sur 1 an. Indicateur de dynamisme commercial.",
    # Dividendes
    "Dividend Yield": "Dividende annuel / Prix de l'action. Rendement en cash pour l'actionnaire. > 3% = généreux.",
    "Dividende/Act":  "Montant brut du dividende versé par action sur 12 mois.",
    "Payout Ratio":   "Part du bénéfice distribué en dividendes. > 80% = peu de marge pour investir.",
    "Ex-Date":        "Date à partir de laquelle tu dois posséder l'action pour recevoir le prochain dividende.",
    # Structure financière
    "Dette/Cap. prop.":"Dette totale / Capitaux propres. Levier financier. > 2 = entreprise très endettée.",
    "Total Cash":     "Liquidités et équivalents de trésorerie disponibles immédiatement.",
    "Total Dette":    "Ensemble des dettes financières de l'entreprise (court et long terme).",
    "Free Cash Flow": "Cash généré après investissements. Argent réellement disponible pour dividendes, rachats ou croissance.",
    "Op. Cash Flow":  "Cash généré par l'activité opérationnelle. Plus fiable que le bénéfice comptable.",
    "EPS (TTM)":      "Bénéfice par action sur 12 mois. Base de calcul du P/E.",
    "EPS Forward":    "Bénéfice par action estimé pour les 12 prochains mois.",
    "Earn. Growth":   "Croissance des bénéfices sur 1 an. Moteur principal de la hausse du cours.",
    # Stats
    "Rend. total":    "Performance totale du prix sur la période sélectionnée (hors dividendes).",
    "Rend. moy. annuel":"Rendement moyen extrapolé sur 1 an depuis la période analysée.",
    "Vol. annualisée":"Volatilité annualisée = écart-type des rendements × √252. Mesure l'amplitude des fluctuations.",
    "Sharpe Ratio":   "Rendement / Risque. Sharpe > 1 = bon. > 2 = excellent. Négatif = perd de l'argent pour du risque.",
    "Sortino Ratio":  "Comme le Sharpe, mais ne pénalise que la volatilité négative (baisse). Plus pertinent pour l'investisseur.",
    "Max Drawdown":   "Perte maximale depuis un sommet jusqu'au creux suivant. Mesure le pire scénario historique.",
    "Jours positifs": "% de séances avec une clôture en hausse. > 52% = tendance haussière structurelle.",
    "Asymétrie":      "Skewness. Positive = rendements extrêmes à la hausse plus fréquents. Négative = risque de crash.",
    "Kurtosis":       "Queue de distribution. > 3 = événements extrêmes plus fréquents que la normale (fat tails).",
    "Meilleure jour": "Meilleure séance sur la période. Rappelle pourquoi rester investi est important.",
    "Pire jour":      "Pire séance sur la période. Illustre le risque de volatilité court terme.",
    # Risque
    "VaR 95% (1j)":   "Value at Risk 95%. Perte maximale attendue sur 1 jour avec 95% de probabilité.",
    "VaR 99% (1j)":   "Value at Risk 99%. Perte maximale attendue sur 1 jour avec 99% de probabilité (scénario extrême).",
    "CVaR 95%":       "Conditional VaR. Perte moyenne dans les 5% pires cas. Encore plus conservateur que la VaR.",
    "CVaR 99%":       "Perte moyenne dans les 1% pires cas. Mesure le risque de queue (tail risk).",
    "Downside Risk":  "Volatilité des rendements négatifs uniquement. Mesure le risque de perte réel.",
}


def metric_with_tooltip(label: str, value: str, tooltip_key: str = None, col=None):
    """Affiche un st.metric avec une info-bulle si disponible."""
    tip = TOOLTIPS.get(tooltip_key or label, "")
    target = col if col is not None else st
    target.metric(label=label, value=value, help=tip if tip else None)

# ══════════════════════════════════════════════════════════════
#  TICKER SEARCH — Recherche via Yahoo Finance
# ══════════════════════════════════════════════════════════════

TYPE_LABELS = {
    "EQUITY":   "Action",
    "ETF":      "ETF",
    "INDEX":    "Indice",
    "MUTUALFUND": "Fonds",
    "CURRENCY": "Devise",
    "FUTURE":   "Future",
    "CRYPTOCURRENCY": "Crypto",
    "OPTION":   "Option",
}

@st.cache_data(ttl=60, show_spinner=False)
def rechercher_tickers(query: str) -> list:
    """
    Cherche des tickers via l'API Yahoo Finance autocomplete.
    Retourne une liste de dicts {symbol, name, type, exchange}.
    """
    try:
        q = urllib.parse.quote(query)
        url = f"https://query1.finance.yahoo.com/v1/finance/search?q={q}&quotesCount=15&newsCount=0&enableFuzzyQuery=true"
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode())
        results = []
        for item in data.get("quotes", []):
            results.append({
                "symbol":   item.get("symbol", ""),
                "name":     item.get("longname") or item.get("shortname", "—"),
                "type":     item.get("quoteType", ""),
                "exchange": item.get("exchDisp", item.get("exchange", "")),
            })
        return results
    except Exception:
        return []

# ─────────────────────────────────────────────
#  PAGE CONFIG
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="Apex Markets",
    page_icon="△",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─────────────────────────────────────────────
#  CSS
# ─────────────────────────────────────────────
st.markdown("""
<style>
    .main-header {
        font-size: 2.4rem; font-weight: 800;
        background: linear-gradient(90deg, #00d4ff, #0066ff);
        -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        margin-bottom: 0.2rem;
    }
    .metric-card {
        background: #1e2130; border-radius: 10px;
        padding: 1rem 1.2rem; border-left: 4px solid #0066ff;
    }
    .section-title {
        font-size: 1.3rem; font-weight: 700; color: #00d4ff;
        border-bottom: 1px solid #2a2d3e; padding-bottom: 0.4rem;
        margin-top: 1.4rem; margin-bottom: 0.8rem;
    }
    .badge-positive { color: #00e676; font-weight: 700; }
    .badge-negative { color: #ff5252; font-weight: 700; }
    .stMetric label { font-size: 0.78rem !important; color: #8b92a5 !important; }
    div[data-testid="stMetricValue"] { font-size: 1.1rem !important; }

    /* ── Onglets — style pro ── */
    .stTabs [data-baseweb="tab-list"] {
        gap: 0;
        background: #151725;
        border-radius: 8px;
        padding: 4px;
        border: 1px solid #2a2d3e;
    }
    .stTabs [data-baseweb="tab"] {
        padding: 10px 20px;
        border-radius: 6px;
        font-size: 0.82rem;
        font-weight: 500;
        color: #8b92a5;
        transition: all 0.15s ease;
        border: none;
        background: transparent;
    }
    .stTabs [data-baseweb="tab"]:hover {
        color: #c5cae9;
        background: rgba(0, 102, 255, 0.08);
    }
    .stTabs [aria-selected="true"] {
        color: #ffffff !important;
        background: linear-gradient(135deg, #0055dd, #0066ff) !important;
        font-weight: 600 !important;
        box-shadow: 0 2px 8px rgba(0, 102, 255, 0.25);
    }
    .stTabs [data-baseweb="tab-highlight"] {
        display: none;
    }
    .stTabs [data-baseweb="tab-border"] {
        display: none;
    }
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
#  SECTION 1 — DATA LOADING
# ══════════════════════════════════════════════════════════════

@st.cache_data(ttl=300, show_spinner=False)
def load_data(symbol: str, period: str, interval: str) -> pd.DataFrame:
    """Télécharge l'historique de prix via yfinance."""
    try:
        ticker = yf.Ticker(symbol)
        df = ticker.history(period=period, interval=interval, auto_adjust=True)
        if df.empty:
            return pd.DataFrame()
        df.index = pd.to_datetime(df.index)
        if df.index.tz is not None:
            df.index = df.index.tz_localize(None)
        return df
    except Exception as e:
        st.error(f"Erreur chargement données : {e}")
        return pd.DataFrame()


@st.cache_data(ttl=600, show_spinner=False)
def get_ticker_info(symbol: str) -> dict:
    """Récupère le dictionnaire info complet du ticker."""
    try:
        ticker = yf.Ticker(symbol)
        info = ticker.info or {}
        # Fallback via fast_info pour les champs manquants
        try:
            fi = ticker.fast_info
            if not info.get("currency") and hasattr(fi, "currency"):
                info["currency"] = fi.currency
            if not info.get("longName") and not info.get("shortName"):
                if hasattr(fi, "long_name") and fi.long_name:
                    info["longName"] = fi.long_name
                elif hasattr(fi, "short_name") and fi.short_name:
                    info["shortName"] = fi.short_name
            if not info.get("marketCap") and hasattr(fi, "market_cap"):
                info["marketCap"] = fi.market_cap
        except Exception:
            pass
        return info
    except Exception:
        return {}


@st.cache_data(ttl=600, show_spinner=False)
def get_fundamentals(symbol: str) -> dict:
    """
    Récupère toutes les données fondamentales disponibles :
    financials, balance sheet, cash flow, dividendes, recommandations…
    """
    result = {}
    ticker = yf.Ticker(symbol)
    for attr in [
        "financials", "quarterly_financials",
        "balance_sheet", "quarterly_balance_sheet",
        "cashflow", "quarterly_cashflow",
        "dividends", "splits",
        "recommendations", "calendar",
        "earnings", "earnings_dates",
        "sustainability",
        "institutional_holders", "major_holders",
        "analyst_price_targets",
    ]:
        try:
            val = getattr(ticker, attr, None)
            if val is not None:
                result[attr] = val
        except Exception:
            pass
    return result


# ══════════════════════════════════════════════════════════════
#  SECTION 2 — TECHNICAL INDICATORS
# ══════════════════════════════════════════════════════════════

def compute_indicators(df: pd.DataFrame) -> pd.DataFrame:
    """
    Calcule l'ensemble des indicateurs techniques :
    SMA, EMA, MACD, RSI, Stochastique, Bollinger Bands, ATR, OBV.
    """
    if df.empty or len(df) < 5:
        return df

    close = df["Close"]
    high  = df["High"]
    low   = df["Low"]
    vol   = df["Volume"]

    # ── Moving Averages ──────────────────────────────────────
    for w in [5, 10, 20, 50, 100, 200]:
        df[f"SMA{w}"] = close.rolling(w).mean()
        df[f"EMA{w}"] = close.ewm(span=w, adjust=False).mean()

    # ── MACD ────────────────────────────────────────────────
    ema12 = close.ewm(span=12, adjust=False).mean()
    ema26 = close.ewm(span=26, adjust=False).mean()
    df["MACD"]        = ema12 - ema26
    df["MACD_Signal"] = df["MACD"].ewm(span=9, adjust=False).mean()
    df["MACD_Hist"]   = df["MACD"] - df["MACD_Signal"]

    # ── RSI ─────────────────────────────────────────────────
    delta = close.diff()
    gain  = delta.clip(lower=0).rolling(14).mean()
    loss  = (-delta.clip(upper=0)).rolling(14).mean()
    rs    = gain / loss.replace(0, np.nan)
    df["RSI"] = 100 - (100 / (1 + rs))

    # ── Stochastic ──────────────────────────────────────────
    low14  = low.rolling(14).min()
    high14 = high.rolling(14).max()
    df["Stoch_K"] = 100 * (close - low14) / (high14 - low14 + 1e-9)
    df["Stoch_D"] = df["Stoch_K"].rolling(3).mean()

    # ── Bollinger Bands ─────────────────────────────────────
    sma20       = close.rolling(20).mean()
    std20       = close.rolling(20).std()
    df["BB_Mid"]   = sma20
    df["BB_Upper"] = sma20 + 2 * std20
    df["BB_Lower"] = sma20 - 2 * std20
    df["BB_Width"] = (df["BB_Upper"] - df["BB_Lower"]) / sma20

    # ── ATR ─────────────────────────────────────────────────
    prev_close   = close.shift(1)
    tr           = pd.concat([
        high - low,
        (high - prev_close).abs(),
        (low  - prev_close).abs(),
    ], axis=1).max(axis=1)
    df["ATR"] = tr.rolling(14).mean()

    # ── OBV ─────────────────────────────────────────────────
    obv  = (np.sign(close.diff()) * vol).fillna(0).cumsum()
    df["OBV"] = obv
    df["OBV_MA"] = df["OBV"].rolling(20).mean()

    # ── Volume MA ───────────────────────────────────────────
    df["Vol_MA20"] = vol.rolling(20).mean()

    # ── Returns & Drawdown ──────────────────────────────────
    df["Daily_Return"]   = close.pct_change()
    df["Log_Return"]     = np.log(close / close.shift(1))
    df["Cum_Return"]     = (1 + df["Daily_Return"]).cumprod() - 1
    roll_max             = close.cummax()
    df["Drawdown"]       = (close - roll_max) / roll_max

    return df


# ══════════════════════════════════════════════════════════════
#  SECTION 3 — STATISTICS
# ══════════════════════════════════════════════════════════════

def compute_statistics(df: pd.DataFrame) -> dict:
    """Calcule les statistiques de performance et risque."""
    stats = {}
    if df.empty or "Daily_Return" not in df.columns:
        return stats

    ret = df["Daily_Return"].dropna()
    if len(ret) < 2:
        return stats

    trading_days = 252
    stats["mean_daily_return"]    = ret.mean()
    stats["mean_annual_return"]   = ret.mean() * trading_days
    stats["volatility_daily"]     = ret.std()
    stats["volatility_annual"]    = ret.std() * np.sqrt(trading_days)
    # Taux sans risque annualisé (US 10Y ~4.3%, ajustable)
    rf_annual = 0.043
    rf_daily  = rf_annual / trading_days
    excess    = ret - rf_daily
    stats["sharpe_ratio"]         = (excess.mean() / ret.std()) * np.sqrt(trading_days) if ret.std() > 0 else 0
    stats["sortino_ratio"]        = (excess.mean() / ret[ret < 0].std()) * np.sqrt(trading_days) if ret[ret < 0].std() > 0 else 0
    stats["max_drawdown"]         = df["Drawdown"].min() if "Drawdown" in df.columns else np.nan
    stats["skewness"]             = ret.skew()
    stats["kurtosis"]             = ret.kurt()
    stats["positive_days_pct"]    = (ret > 0).mean() * 100
    stats["best_day"]             = ret.max()
    stats["worst_day"]            = ret.min()
    stats["total_return"]         = (df["Close"].iloc[-1] / df["Close"].iloc[0] - 1) if len(df) > 1 else 0
    return stats


def compute_risk_metrics(df: pd.DataFrame, confidence: float = 0.95) -> dict:
    """Calcule les métriques de risque : VaR, CVaR, downside risk."""
    risk = {}
    if df.empty or "Daily_Return" not in df.columns:
        return risk

    ret = df["Daily_Return"].dropna()
    if len(ret) < 10:
        return risk

    # VaR historique
    risk["VaR_95"]   = np.percentile(ret, 5)
    risk["VaR_99"]   = np.percentile(ret, 1)
    # CVaR (Expected Shortfall)
    risk["CVaR_95"]  = ret[ret <= risk["VaR_95"]].mean()
    risk["CVaR_99"]  = ret[ret <= risk["VaR_99"]].mean()
    # Downside deviation
    risk["downside_dev"] = ret[ret < 0].std() * np.sqrt(252)
    # Annualised vol
    risk["annual_vol"]   = ret.std() * np.sqrt(252)
    return risk


# ══════════════════════════════════════════════════════════════
#  SECTION 4 — CHARTS
# ══════════════════════════════════════════════════════════════

CHART_BG   = "#0e1117"
GRID_COLOR = "#2a2d3e"
UP_COLOR   = "#00e676"
DOWN_COLOR = "#ff5252"
ACCENT     = "#00aaff"


def create_price_chart(df: pd.DataFrame, symbol: str, options: dict) -> go.Figure:
    """
    Graphique principal : candlestick + volume + overlays configurables.
    """
    rows    = 2 if options.get("show_volume", True) else 1
    heights = [0.72, 0.28] if rows == 2 else [1]

    fig = make_subplots(
        rows=rows, cols=1,
        shared_xaxes=True,
        vertical_spacing=0.02,
        row_heights=heights,
    )

    # ── Candlestick ─────────────────────────────────────────
    fig.add_trace(go.Candlestick(
        x=df.index, open=df["Open"], high=df["High"],
        low=df["Low"], close=df["Close"],
        name="Prix",
        increasing_line_color=UP_COLOR,
        decreasing_line_color=DOWN_COLOR,
        increasing_fillcolor=UP_COLOR,
        decreasing_fillcolor=DOWN_COLOR,
    ), row=1, col=1)

    # ── Bollinger Bands ─────────────────────────────────────
    if options.get("bollinger") and "BB_Upper" in df.columns:
        for col, dash, name in [
            ("BB_Upper", "dash", "BB Upper"),
            ("BB_Mid",   "dot",  "BB Mid"),
            ("BB_Lower", "dash", "BB Lower"),
        ]:
            fig.add_trace(go.Scatter(
                x=df.index, y=df[col], name=name,
                line=dict(color="#ff9800", dash=dash, width=1),
                opacity=0.75,
            ), row=1, col=1)
        fig.add_trace(go.Scatter(
            x=list(df.index) + list(df.index[::-1]),
            y=list(df["BB_Upper"]) + list(df["BB_Lower"][::-1]),
            fill="toself", fillcolor="rgba(255,152,0,0.06)",
            line=dict(color="rgba(0,0,0,0)"), name="BB Band",
        ), row=1, col=1)

    # ── Moving Averages ──────────────────────────────────────
    ma_colors = {"MA20": "#f9a825", "MA50": "#42a5f5",
                 "MA100": "#ab47bc", "MA200": "#ef5350"}
    for ma_key, color in ma_colors.items():
        if options.get(ma_key):
            num  = ma_key.replace("MA", "")
            col  = f"SMA{num}"
            if col in df.columns:
                fig.add_trace(go.Scatter(
                    x=df.index, y=df[col], name=ma_key,
                    line=dict(color=color, width=1.4),
                ), row=1, col=1)

    # ── Volume ───────────────────────────────────────────────
    if options.get("show_volume", True) and "Volume" in df.columns:
        colors = [UP_COLOR if df["Close"].iloc[i] >= df["Open"].iloc[i] else DOWN_COLOR
                  for i in range(len(df))]
        fig.add_trace(go.Bar(
            x=df.index, y=df["Volume"], name="Volume",
            marker_color=colors, opacity=0.65,
        ), row=2, col=1)
        if "Vol_MA20" in df.columns:
            fig.add_trace(go.Scatter(
                x=df.index, y=df["Vol_MA20"], name="Vol MA20",
                line=dict(color="#ff9800", width=1.2),
            ), row=2, col=1)

    # ── Axe X catégoriel (comme TradingView) : chaque bougie est
    #    espacée de manière égale, aucun gap weekend/nuit ──
    # Formater les labels selon l'intervalle
    itv = options.get("interval", "1d")
    if itv in ("1m", "5m", "15m", "30m"):
        tick_fmt = "%d %b %Hh%M"
    elif itv == "1h":
        tick_fmt = "%d %b %Hh"
    elif itv in ("1d",):
        tick_fmt = "%d %b %Y"
    elif itv == "1wk":
        tick_fmt = "%d %b %Y"
    else:
        tick_fmt = "%b %Y"

    # Réduire le nombre de ticks pour la lisibilité
    n = len(df)
    max_ticks = 30
    tick_step = max(1, n // max_ticks)
    tick_vals = list(df.index[::tick_step])
    tick_text = [t.strftime(tick_fmt) for t in tick_vals]

    fig.update_layout(
        title=dict(text=f"<b>{symbol}</b> — Analyse du Prix", font_size=16),
        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
        xaxis_rangeslider_visible=False,
        legend=dict(bgcolor="rgba(0,0,0,0.4)", font_color="#ccc"),
        hovermode="x unified",
        height=560,
        margin=dict(l=10, r=10, t=40, b=10),
        dragmode="pan",
    )
    # Appliquer l'axe catégoriel sur toutes les lignes
    for r in range(1, rows + 1):
        fig.update_xaxes(
            type="category",
            tickvals=tick_vals,
            ticktext=tick_text if r == rows else [""] * len(tick_vals),
            tickangle=-45,
            row=r, col=1,
        )

    if options.get("log_scale"):
        fig.update_yaxes(type="log", row=1, col=1)

    for ax in ["xaxis", "xaxis2", "yaxis", "yaxis2"]:
        fig.update_layout(**{ax: dict(
            gridcolor=GRID_COLOR, zerolinecolor=GRID_COLOR,
            tickfont_color="#8b92a5",
        )})
    return fig


def create_technical_charts(df: pd.DataFrame, interval: str = "1d") -> go.Figure:
    """Panneau RSI + MACD + Stochastique + OBV."""
    fig = make_subplots(
        rows=4, cols=1,
        shared_xaxes=True, vertical_spacing=0.03,
        subplot_titles=("RSI (14)", "MACD (12/26/9)",
                        "Stochastique (14,3)", "OBV"),
        row_heights=[0.25, 0.3, 0.25, 0.2],
    )

    # ── RSI ─────────────────────────────────────────────────
    if "RSI" in df.columns:
        fig.add_trace(go.Scatter(x=df.index, y=df["RSI"], name="RSI",
                                 line=dict(color=ACCENT, width=1.5)), row=1, col=1)
        for lvl, color, label in [(70, DOWN_COLOR, "Surachat"),
                                  (30, UP_COLOR,   "Survente")]:
            fig.add_hline(y=lvl, line_color=color, line_dash="dash",
                          opacity=0.6, row=1, col=1)
        fig.add_hrect(y0=30, y1=70, fillcolor="rgba(255,255,255,0.03)",
                      line_width=0, row=1, col=1)

    # ── MACD ────────────────────────────────────────────────
    if "MACD" in df.columns:
        fig.add_trace(go.Scatter(x=df.index, y=df["MACD"], name="MACD",
                                 line=dict(color="#42a5f5", width=1.5)), row=2, col=1)
        fig.add_trace(go.Scatter(x=df.index, y=df["MACD_Signal"], name="Signal",
                                 line=dict(color="#f9a825", width=1.5)), row=2, col=1)
        colors = [UP_COLOR if v >= 0 else DOWN_COLOR for v in df["MACD_Hist"].fillna(0)]
        fig.add_trace(go.Bar(x=df.index, y=df["MACD_Hist"], name="Hist",
                             marker_color=colors, opacity=0.7), row=2, col=1)

    # ── Stochastique ─────────────────────────────────────────
    if "Stoch_K" in df.columns:
        fig.add_trace(go.Scatter(x=df.index, y=df["Stoch_K"], name="%K",
                                 line=dict(color="#ce93d8", width=1.4)), row=3, col=1)
        fig.add_trace(go.Scatter(x=df.index, y=df["Stoch_D"], name="%D",
                                 line=dict(color="#f9a825", width=1.2,
                                           dash="dot")), row=3, col=1)
        for lvl in [80, 20]:
            fig.add_hline(y=lvl, line_color="#888", line_dash="dash",
                          opacity=0.5, row=3, col=1)

    # ── OBV ─────────────────────────────────────────────────
    if "OBV" in df.columns:
        fig.add_trace(go.Scatter(x=df.index, y=df["OBV"], name="OBV",
                                 line=dict(color=UP_COLOR, width=1.4)), row=4, col=1)
        if "OBV_MA" in df.columns:
            fig.add_trace(go.Scatter(x=df.index, y=df["OBV_MA"], name="OBV MA20",
                                     line=dict(color="#f9a825", width=1,
                                               dash="dot")), row=4, col=1)

    # Axe catégoriel : pas de gaps
    n = len(df)
    tick_step = max(1, n // 25)
    tick_vals = list(df.index[::tick_step])
    if interval in ("1m", "5m", "15m", "30m"):
        tick_text = [t.strftime("%d %b %Hh%M") for t in tick_vals]
    elif interval == "1h":
        tick_text = [t.strftime("%d %b %Hh") for t in tick_vals]
    else:
        tick_text = [t.strftime("%d %b %Y") for t in tick_vals]

    fig.update_layout(
        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
        height=680, hovermode="x unified",
        showlegend=True,
        legend=dict(bgcolor="rgba(0,0,0,0.4)", font_color="#ccc"),
        margin=dict(l=10, r=10, t=30, b=10),
    )
    for i in range(1, 5):
        fig.update_xaxes(
            type="category",
            tickvals=tick_vals,
            ticktext=tick_text if i == 4 else [""] * len(tick_vals),
            tickangle=-45,
            row=i, col=1,
        )
    for ax in ["xaxis", "xaxis2", "xaxis3", "xaxis4",
               "yaxis", "yaxis2", "yaxis3", "yaxis4"]:
        fig.update_layout(**{ax: dict(
            gridcolor=GRID_COLOR, tickfont_color="#8b92a5",
        )})
    return fig


def create_drawdown_chart(df: pd.DataFrame, symbol: str, interval: str = "1d") -> go.Figure:
    """Drawdown historique."""
    fig = go.Figure()
    if "Drawdown" in df.columns:
        fig.add_trace(go.Scatter(
            x=df.index, y=df["Drawdown"] * 100,
            fill="tozeroy", name="Drawdown",
            line=dict(color=DOWN_COLOR, width=1.5),
            fillcolor="rgba(255,82,82,0.2)",
        ))
    n = len(df)
    tick_step = max(1, n // 25)
    tick_vals = list(df.index[::tick_step])
    if interval in ("1m", "5m", "15m", "30m"):
        tick_text = [t.strftime("%d %b %Hh%M") for t in tick_vals]
    elif interval == "1h":
        tick_text = [t.strftime("%d %b %Hh") for t in tick_vals]
    else:
        tick_text = [t.strftime("%d %b %Y") for t in tick_vals]
    fig.update_layout(
        title=f"<b>{symbol}</b> — Drawdown historique (%)",
        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
        yaxis=dict(ticksuffix="%", gridcolor=GRID_COLOR, tickfont_color="#8b92a5"),
        xaxis=dict(type="category", tickvals=tick_vals, ticktext=tick_text,
                   tickangle=-45, gridcolor=GRID_COLOR, tickfont_color="#8b92a5"),
        height=300, margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


def create_returns_distribution(df: pd.DataFrame) -> go.Figure:
    """Distribution des rendements journaliers avec courbe loi normale, skewness et kurtosis."""
    fig = go.Figure()
    if "Daily_Return" not in df.columns:
        return fig
    ret = df["Daily_Return"].dropna() * 100
    mean_r = ret.mean()
    std_r  = ret.std()
    skew_r = ret.skew()
    kurt_r = ret.kurt()

    # Histogramme (densité pour superposition avec la courbe)
    fig.add_trace(go.Histogram(
        x=ret, nbinsx=60, name="Rendements",
        marker_color=ACCENT, opacity=0.75,
        histnorm="probability density",
    ))

    # Courbe loi normale ajustée
    x_range = np.linspace(ret.min(), ret.max(), 300)
    y_normal = scipy_stats.norm.pdf(x_range, loc=mean_r, scale=std_r)
    fig.add_trace(go.Scatter(
        x=x_range, y=y_normal,
        name="Loi Normale",
        line=dict(color="#ff9800", width=2.2, dash="solid"),
    ))

    # Lignes verticales : moyenne, ±1σ, ±2σ
    fig.add_vline(x=mean_r, line_color="#f9a825", line_dash="dash",
                  annotation_text=f"μ: {mean_r:.2f}%",
                  annotation_font_color="#f9a825", annotation_position="top right")
    for k, alpha in [(1, 0.7), (2, 0.4)]:
        for sign in [-1, 1]:
            x_val = mean_r + sign * k * std_r
            fig.add_vline(x=x_val, line_color="#ab47bc",
                          line_dash="dot", opacity=alpha,
                          annotation_text=f"{sign*k:+d}σ",
                          annotation_font_color="#ab47bc")

    # Annotation skewness / kurtosis
    skew_label = "→ Queue droite (gains extremes)" if skew_r > 0.3 else ("→ Queue gauche (pertes extremes)" if skew_r < -0.3 else "→ Symétrique")
    kurt_label = "→ Fat tails (événements extrêmes fréquents)" if kurt_r > 1 else ("→ Queue fine" if kurt_r < -0.5 else "→ Distribution normale")
    annotation_text = (
        f"<b>Skewness :</b> {skew_r:.3f}  {skew_label}<br>"
        f"<b>Kurtosis :</b> {kurt_r:.3f}  {kurt_label}"
    )
    fig.add_annotation(
        x=0.01, y=0.98, xref="paper", yref="paper",
        text=annotation_text,
        showarrow=False, align="left",
        bgcolor="rgba(30,33,48,0.85)",
        bordercolor="#2a2d3e", borderwidth=1,
        font=dict(color="#c5cae9", size=11),
        xanchor="left", yanchor="top",
    )

    fig.update_layout(
        title="Distribution des rendements journaliers (%) — vs Loi Normale",
        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
        xaxis=dict(gridcolor=GRID_COLOR, ticksuffix="%", tickfont_color="#8b92a5"),
        yaxis=dict(gridcolor=GRID_COLOR, tickfont_color="#8b92a5", title="Densité"),
        legend=dict(bgcolor="rgba(0,0,0,0.4)", font_color="#ccc"),
        height=480, margin=dict(l=50, r=30, t=50, b=40),
        bargap=0.06,
    )
    return fig


def create_cumulative_returns_chart(dfs: dict) -> go.Figure:
    """Rendements cumulés de plusieurs tickers."""
    fig = go.Figure()
    colors = ["#00aaff", "#00e676", "#f9a825", "#ef5350",
              "#ab47bc", "#ff9800", "#26c6da"]
    for i, (sym, df) in enumerate(dfs.items()):
        if df.empty or "Cum_Return" not in df.columns:
            continue
        color = colors[i % len(colors)]
        fig.add_trace(go.Scatter(
            x=df.index, y=df["Cum_Return"] * 100,
            name=sym, line=dict(color=color, width=1.8),
        ))
    fig.add_hline(y=0, line_color="#666", line_dash="dash")
    fig.update_layout(
        title="Rendements cumulés comparés (%)",
        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
        yaxis=dict(ticksuffix="%", gridcolor=GRID_COLOR, tickfont_color="#8b92a5"),
        xaxis=dict(gridcolor=GRID_COLOR, tickfont_color="#8b92a5"),
        height=400, hovermode="x unified",
        legend=dict(bgcolor="rgba(0,0,0,0.4)", font_color="#ccc"),
        margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


def create_correlation_heatmap(dfs: dict) -> go.Figure:
    """Heatmap de corrélation des rendements."""
    rets = {}
    for sym, df in dfs.items():
        if not df.empty and "Daily_Return" in df.columns:
            rets[sym] = df["Daily_Return"]
    if len(rets) < 2:
        return go.Figure()

    corr = pd.DataFrame(rets).corr()
    fig = go.Figure(go.Heatmap(
        z=corr.values, x=corr.columns, y=corr.index,
        colorscale="RdBu", zmin=-1, zmax=1,
        text=corr.round(2).values, texttemplate="%{text}",
        colorbar=dict(tickfont_color="#ccc"),
    ))
    fig.update_layout(
        title="Matrice de corrélation des rendements",
        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
        xaxis=dict(tickfont_color="#8b92a5"),
        yaxis=dict(tickfont_color="#8b92a5"),
        height=400, margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


def create_monthly_returns_heatmap(df: pd.DataFrame, symbol: str) -> go.Figure:
    """Heatmap des rendements mensuels par année."""
    if df.empty or "Daily_Return" not in df.columns:
        return go.Figure()

    monthly = df["Daily_Return"].resample("ME").apply(
        lambda x: (1 + x).prod() - 1
    ) * 100
    monthly = monthly.to_frame("Return")
    monthly["Year"]  = monthly.index.year
    monthly["Month"] = monthly.index.month

    pivot = monthly.pivot_table(index="Year", columns="Month", values="Return")
    month_labels = ["Jan","Fév","Mar","Avr","Mai","Jun",
                    "Jul","Aoû","Sep","Oct","Nov","Déc"]
    pivot.columns = [month_labels[m - 1] for m in pivot.columns]

    # Préparer les annotations : chaîne vide pour les NaN, sinon "X.X%"
    text_display = []
    for row in pivot.values:
        text_row = []
        for v in row:
            text_row.append(f"{v:.1f}%" if not np.isnan(v) else "")
        text_display.append(text_row)

    fig = go.Figure(go.Heatmap(
        z=pivot.values, x=pivot.columns, y=pivot.index.astype(str),
        colorscale="RdYlGn", zmid=0,
        text=text_display,
        texttemplate="%{text}",
        hovertemplate="Mois: %{x}<br>Année: %{y}<br>Rendement: %{z:.1f}%<extra></extra>",
        colorbar=dict(tickfont_color="#ccc", ticksuffix="%"),
    ))
    fig.update_layout(
        title=f"<b>{symbol}</b> — Rendements mensuels (%)",
        paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
        xaxis=dict(tickfont_color="#8b92a5"),
        yaxis=dict(tickfont_color="#8b92a5"),
        height=max(250, len(pivot) * 32 + 80),
        margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


# ══════════════════════════════════════════════════════════════
#  SECTION 5 — DISPLAY HELPERS
# ══════════════════════════════════════════════════════════════

def fmt(value, decimals=2, suffix="", prefix="", na_str="N/A") -> str:
    """Formate une valeur numérique en format français (espace milliers, virgule décimale)."""
    try:
        if value is None or (isinstance(value, float) and np.isnan(value)):
            return na_str
        formatted = f"{value:,.{decimals}f}"
        # Conversion en format FR : 1,234.56 → 1 234,56
        formatted = formatted.replace(",", "\u00a0").replace(".", ",")
        return f"{prefix}{formatted}{suffix}"
    except Exception:
        return str(value)


def fmt_large(value, currency: str = "") -> str:
    """Formate les grandes valeurs (milliards / millions)."""
    sym = {"USD": "$", "EUR": "€", "GBP": "£", "JPY": "¥", "CHF": "CHF "}.get(currency, currency + " " if currency else "")
    try:
        v = float(value)
        if abs(v) >= 1e12:
            return f"{v/1e12:.2f} T{sym}"
        if abs(v) >= 1e9:
            return f"{v/1e9:.2f} Md{sym}"
        if abs(v) >= 1e6:
            return f"{v/1e6:.2f} M{sym}"
        return fmt(v)
    except Exception:
        return "N/A"


def display_fundamentals(info: dict):
    """Affiche l'ensemble des données fondamentales."""
    if not info:
        st.warning("Données fondamentales indisponibles.")
        return

    cur = info.get("currency", "USD")
    cur_sym = {"USD": "$", "EUR": "€", "GBP": "£", "JPY": "¥", "CHF": "CHF "}.get(cur, cur + " ")

    # ─── Section Prix ────────────────────────────────────────
    st.markdown('<p class="section-title">Données de Prix</p>', unsafe_allow_html=True)
    cols = st.columns(4)
    price_fields = [
        ("Prix actuel",     "currentPrice"),
        ("Open",            "open"),
        ("High (jour)",     "dayHigh"),
        ("Low (jour)",      "dayLow"),
        ("Clôture préc.",   "previousClose"),
        ("52-sem. High",    "fiftyTwoWeekHigh"),
        ("52-sem. Low",     "fiftyTwoWeekLow"),
        ("Volume",          "volume"),
    ]
    for i, (label, key) in enumerate(price_fields):
        val = info.get(key)
        with cols[i % 4]:
            if key == "volume":
                metric_with_tooltip(label, fmt_large(val) if val else "N/A")
            else:
                metric_with_tooltip(label, fmt(val, suffix="", prefix=cur_sym) if val else "N/A")

    # ─── Market Cap & Beta ───────────────────────────────────
    st.markdown('<p class="section-title">Capitalisation & Valorisation</p>', unsafe_allow_html=True)
    cols = st.columns(4)
    val_fields = [
        ("Market Cap",    "marketCap",         lambda v: fmt_large(v, cur)),
        ("Enterprise V.", "enterpriseValue",   lambda v: fmt_large(v, cur)),
        ("Beta",          "beta",              lambda v: fmt(v)),
        ("Float",         "floatShares",       fmt_large),
        ("P/E (TTM)",     "trailingPE",        lambda v: fmt(v, 1)),
        ("P/E Forward",   "forwardPE",         lambda v: fmt(v, 1)),
        ("PEG",           "pegRatio",          lambda v: fmt(v, 2)),
        ("P/B",           "priceToBook",       lambda v: fmt(v, 2)),
        ("P/S",           "priceToSalesTrailing12Months", lambda v: fmt(v, 2)),
        ("EV/EBITDA",     "enterpriseToEbitda",lambda v: fmt(v, 1)),
        ("EV/Rev.",       "enterpriseToRevenue",lambda v: fmt(v, 2)),
        ("52w perf",      "52WeekChange",      lambda v: fmt(v*100, 1, "%")),
    ]
    for i, (label, key, fn) in enumerate(val_fields):
        val = info.get(key)
        with cols[i % 4]:
            metric_with_tooltip(label, fn(val) if val is not None else "N/A")

    # ─── Rentabilité ─────────────────────────────────────────
    st.markdown('<p class="section-title">Rentabilité</p>', unsafe_allow_html=True)
    cols = st.columns(4)
    profit_fields = [
        ("ROE",             "returnOnEquity",    lambda v: fmt(v*100, 1, "%")),
        ("ROA",             "returnOnAssets",    lambda v: fmt(v*100, 1, "%")),
        ("Marge brute",     "grossMargins",      lambda v: fmt(v*100, 1, "%")),
        ("Marge opérat.",   "operatingMargins",  lambda v: fmt(v*100, 1, "%")),
        ("Marge nette",     "profitMargins",     lambda v: fmt(v*100, 1, "%")),
        ("EBITDA",          "ebitda",            lambda v: fmt_large(v, cur)),
        ("Chiffre d'aff.",  "totalRevenue",      lambda v: fmt_large(v, cur)),
        ("Rev. Growth",     "revenueGrowth",     lambda v: fmt(v*100, 1, "%")),
    ]
    for i, (label, key, fn) in enumerate(profit_fields):
        val = info.get(key)
        with cols[i % 4]:
            metric_with_tooltip(label, fn(val) if val is not None else "N/A")

    # ─── Dividendes ──────────────────────────────────────────
    st.markdown('<p class="section-title">Dividendes</p>', unsafe_allow_html=True)
    cols = st.columns(4)
    div_fields = [
        ("Dividend Yield",  "dividendYield",   lambda v: fmt(v, 2, "%")),
        ("Dividende/Act",   "dividendRate",    lambda v: fmt(v, 2, cur_sym)),
        ("Payout Ratio",    "payoutRatio",     lambda v: fmt(v*100, 1, "%")),
        ("Ex-Date",         "exDividendDate",  lambda v: str(v) if v else "N/A"),
    ]
    for i, (label, key, fn) in enumerate(div_fields):
        val = info.get(key)
        with cols[i % 4]:
            metric_with_tooltip(label, fn(val) if val is not None else "N/A")

    # ─── Structure financière ────────────────────────────────
    st.markdown('<p class="section-title">Structure Financière</p>', unsafe_allow_html=True)
    cols = st.columns(4)
    struct_fields = [
        ("Dette/Cap. prop.", "debtToEquity",      lambda v: fmt(v, 2)),
        ("Total Cash",       "totalCash",          lambda v: fmt_large(v, cur)),
        ("Total Dette",      "totalDebt",          lambda v: fmt_large(v, cur)),
        ("Free Cash Flow",   "freeCashflow",       lambda v: fmt_large(v, cur)),
        ("Op. Cash Flow",    "operatingCashflow",  lambda v: fmt_large(v, cur)),
        ("EPS (TTM)",        "trailingEps",        lambda v: fmt(v, 2, "$")),
        ("EPS Forward",      "forwardEps",         lambda v: fmt(v, 2, "$")),
        ("Earn. Growth",     "earningsGrowth",     lambda v: fmt(v*100, 1, "%")),
    ]
    for i, (label, key, fn) in enumerate(struct_fields):
        val = info.get(key)
        with cols[i % 4]:
            metric_with_tooltip(label, fn(val) if val is not None else "N/A")


def display_technical_analysis(df: pd.DataFrame):
    """Signaux de trading basés sur les indicateurs."""
    if df.empty:
        return

    signals = []
    last = df.iloc[-1]

    # RSI Signal
    if "RSI" in df.columns and not np.isnan(last["RSI"]):
        rsi_val = last["RSI"]
        if rsi_val > 70:
            signals.append(("RSI", f"{rsi_val:.1f}", "Surachat", "negative"))
        elif rsi_val < 30:
            signals.append(("RSI", f"{rsi_val:.1f}", "Survente", "positive"))
        else:
            signals.append(("RSI", f"{rsi_val:.1f}", "Neutre", "neutral"))

    # MA Crossover Signal
    if "SMA20" in df.columns and "SMA50" in df.columns:
        if len(df) >= 2:
            prev = df.iloc[-2]
            cross_up   = last["SMA20"] > last["SMA50"]  and prev["SMA20"] <= prev["SMA50"]
            cross_down = last["SMA20"] < last["SMA50"]  and prev["SMA20"] >= prev["SMA50"]
            pos_cross  = last["SMA20"] > last["SMA50"]
            if cross_up:
                signals.append(("MA20/MA50", "↑ Croisement", "Golden Cross", "positive"))
            elif cross_down:
                signals.append(("MA20/MA50", "↓ Croisement", "Death Cross", "negative"))
            elif pos_cross:
                signals.append(("MA20/MA50", "MA20 > MA50", "Tendance haussière", "positive"))
            else:
                signals.append(("MA20/MA50", "MA20 < MA50", "Tendance baissière", "negative"))

    # MACD Signal
    if "MACD" in df.columns and "MACD_Signal" in df.columns:
        if len(df) >= 2:
            prev = df.iloc[-2]
            bull_cross = (last["MACD"] > last["MACD_Signal"]
                          and prev["MACD"] <= prev["MACD_Signal"])
            bear_cross = (last["MACD"] < last["MACD_Signal"]
                          and prev["MACD"] >= prev["MACD_Signal"])
            if bull_cross:
                signals.append(("MACD", "Croisement haussier", "Signal d'achat", "positive"))
            elif bear_cross:
                signals.append(("MACD", "Croisement baissier", "Signal de vente", "negative"))
            elif last["MACD"] > last["MACD_Signal"]:
                signals.append(("MACD", "MACD > Signal", "Momentum positif", "positive"))
            else:
                signals.append(("MACD", "MACD < Signal", "Momentum négatif", "negative"))

    # Bollinger Band Position
    if "BB_Upper" in df.columns:
        price = last["Close"]
        if price > last["BB_Upper"]:
            signals.append(("Bollinger", f"Prix: {price:.2f}", "Au-dessus BB supér.", "negative"))
        elif price < last["BB_Lower"]:
            signals.append(("Bollinger", f"Prix: {price:.2f}", "En-dessous BB infér.", "positive"))
        else:
            pct = (price - last["BB_Lower"]) / (last["BB_Upper"] - last["BB_Lower"]) * 100
            signals.append(("Bollinger", f"{pct:.0f}% dans la bande", "Intérieur BB", "neutral"))

    # Display signals
    st.markdown('<p class="section-title">Signaux de Trading</p>', unsafe_allow_html=True)
    cols = st.columns(len(signals) if signals else 1)
    for i, (ind, val, sig, kind) in enumerate(signals):
        with cols[i]:
            color = "#00e676" if kind == "positive" else "#ff5252" if kind == "negative" else "#aaa"
            st.markdown(f"""
            <div style="background:#1e2130;border-radius:8px;padding:0.8rem 1rem;
                        border-left:3px solid {color};">
                <div style="color:#8b92a5;font-size:0.75rem;">{ind}</div>
                <div style="font-size:0.95rem;font-weight:600;color:#fff;">{val}</div>
                <div style="font-size:0.8rem;color:{color};">{sig}</div>
            </div>""", unsafe_allow_html=True)


def display_statistics(stats: dict):
    """Affiche les statistiques de performance."""
    if not stats:
        return

    st.markdown('<p class="section-title">Performance</p>', unsafe_allow_html=True)
    cols = st.columns(4)
    fields = [
        ("Rend. total",       "total_return",       lambda v: fmt(v*100, 2, "%")),
        ("Rend. moy. annuel", "mean_annual_return",  lambda v: fmt(v*100, 2, "%")),
        ("Vol. annualisée",   "volatility_annual",   lambda v: fmt(v*100, 2, "%")),
        ("Sharpe Ratio",      "sharpe_ratio",        lambda v: fmt(v, 3)),
        ("Sortino Ratio",     "sortino_ratio",       lambda v: fmt(v, 3)),
        ("Max Drawdown",      "max_drawdown",        lambda v: fmt(v*100, 2, "%")),
        ("Jours positifs",    "positive_days_pct",   lambda v: fmt(v, 1, "%")),
        ("Asymétrie",         "skewness",            lambda v: fmt(v, 3)),
        ("Kurtosis",          "kurtosis",            lambda v: fmt(v, 3)),
        ("Meilleure jour",    "best_day",            lambda v: fmt(v*100, 2, "%")),
        ("Pire jour",         "worst_day",           lambda v: fmt(v*100, 2, "%")),
        ("Rend. vol. quot.",  "volatility_daily",    lambda v: fmt(v*100, 3, "%")),
    ]
    for i, (label, key, fn) in enumerate(fields):
        val = stats.get(key)
        with cols[i % 4]:
            metric_with_tooltip(label, fn(val) if val is not None else "N/A")


def display_risk_analysis(risk: dict):
    """Affiche les métriques de risque."""
    if not risk:
        return

    st.markdown('<p class="section-title">Métriques de Risque</p>', unsafe_allow_html=True)
    cols = st.columns(4)
    fields = [
        ("VaR 95% (1j)",    "VaR_95",     lambda v: fmt(v*100, 2, "%")),
        ("VaR 99% (1j)",    "VaR_99",     lambda v: fmt(v*100, 2, "%")),
        ("CVaR 95%",        "CVaR_95",    lambda v: fmt(v*100, 2, "%")),
        ("CVaR 99%",        "CVaR_99",    lambda v: fmt(v*100, 2, "%")),
        ("Vol. annualisée", "annual_vol", lambda v: fmt(v*100, 2, "%")),
        ("Downside Risk",   "downside_dev",lambda v: fmt(v*100, 2, "%")),
    ]
    for i, (label, key, fn) in enumerate(fields):
        val = risk.get(key)
        with cols[i % 4]:
            metric_with_tooltip(label, fn(val) if val is not None else "N/A")

    # Mini interprétation
    var95 = risk.get("VaR_95", 0)
    if var95:
        st.info(
            f"Avec un niveau de confiance de 95%, la perte maximale attendue sur "
            f"une journée est de **{abs(var95)*100:.2f}%**."
        )


# ══════════════════════════════════════════════════════════════
#  SECTION 6 — FEAR & GREED (proxy simplifié)
# ══════════════════════════════════════════════════════════════

@st.cache_data(ttl=600, show_spinner=False)
def fetch_fear_greed_cnn() -> int:
    """
    Récupère l'indicateur Fear & Greed officiel de CNN Business.
    Fallback à 50 (neutre) en cas d'erreur.
    """
    try:
        url = "https://production.dataviz.cnn.io/index/fearandgreed/graphdata"
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode())
        return int(round(data["fear_and_greed"]["score"]))
    except Exception:
        return 50


def render_fear_greed(score: int):
    """Affiche la jauge Fear & Greed."""
    if score < 25:
        label, color = "Peur extrême", "#ff5252"
    elif score < 45:
        label, color = "Peur", "#ff9800"
    elif score < 55:
        label, color = "Neutre", "#f9a825"
    elif score < 75:
        label, color = "Avidité", "#69f0ae"
    else:
        label, color = "Avidité extrême", "#00e676"

    st.markdown(f"""
    <div style="background:#1e2130;border-radius:8px;padding:0.6rem 0.8rem;text-align:center;">
        <div style="font-size:0.7rem;color:#8b92a5;margin-bottom:0.15rem;letter-spacing:0.03em;">
            Fear & Greed (CNN)
        </div>
        <div style="display:flex;align-items:baseline;justify-content:center;gap:6px;">
            <span style="font-size:1.6rem;font-weight:800;color:{color};line-height:1;">{score}</span>
            <span style="font-size:0.7rem;color:{color};font-weight:600;">{label}</span>
        </div>
        <div style="background:#2a2d3e;border-radius:20px;height:5px;margin-top:0.35rem;">
            <div style="background:{color};width:{score}%;height:5px;
                        border-radius:20px;transition:width 0.5s;"></div>
        </div>
    </div>""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
#  MAIN APP
# ══════════════════════════════════════════════════════════════

def main():
    # ─── Header ──────────────────────────────────────────────
    st.markdown('''
    <div style="display:flex;align-items:center;gap:14px;margin-bottom:0.2rem;">
        <svg width="44" height="44" viewBox="0 0 44 44" fill="none" xmlns="http://www.w3.org/2000/svg">
            <defs>
                <linearGradient id="apex-grad" x1="0" y1="44" x2="44" y2="0">
                    <stop offset="0%" stop-color="#0066ff"/>
                    <stop offset="100%" stop-color="#00d4ff"/>
                </linearGradient>
            </defs>
            <polygon points="22,4 40,38 4,38" fill="none" stroke="url(#apex-grad)" stroke-width="2.8" stroke-linejoin="round"/>
            <polyline points="10,30 18,22 24,26 34,14" stroke="url(#apex-grad)" stroke-width="2.4" stroke-linecap="round" stroke-linejoin="round" fill="none"/>
            <circle cx="34" cy="14" r="2.2" fill="#00d4ff"/>
        </svg>
        <h1 class="main-header" style="margin:0;">Apex Markets</h1>
    </div>''', unsafe_allow_html=True)
    st.markdown("*Plateforme d'analyse financière avancée — investisseurs & traders*")

    # ─── Sidebar ─────────────────────────────────────────────
    with st.sidebar:
        st.markdown("## Configuration")

        symbol_input = st.text_input(
            "Ticker principal",
            value="AAPL",
            help="Ex : AAPL, TSLA, BNP.PA, ^FCHI, ^GSPC",
        ).upper().strip()

        # ── Périodes (labels FR → codes yfinance) ──
        PERIOD_OPTIONS = [
            ("1 semaine",  "5d"),
            ("1 mois",     "1mo"),
            ("3 mois",     "3mo"),
            ("6 mois",     "6mo"),
            ("1 an",       "1y"),
            ("2 ans",      "2y"),
            ("5 ans",      "5y"),
            ("10 ans",     "10y"),
            ("Depuis le 1er janv.", "ytd"),
            ("Maximum",    "max"),
        ]
        period_labels = [p[0] for p in PERIOD_OPTIONS]
        period_codes  = {p[0]: p[1] for p in PERIOD_OPTIONS}
        period_label  = st.selectbox(
            "Période", period_labels, index=4,
            help="Fenêtre temporelle d'analyse. Les intervalles disponibles s'adaptent automatiquement.",
        )
        period = period_codes[period_label]

        # ── Intervalles (labels FR → codes yfinance) ──
        INTERVAL_LABELS = {
            "1m":  "1 minute",   "5m":  "5 minutes",
            "15m": "15 minutes", "30m": "30 minutes",
            "1h":  "1 heure",    "1d":  "1 jour",
            "1wk": "1 semaine",  "1mo": "1 mois",
        }
        interval_map = {
            "5d":   ["1m", "5m", "15m", "30m", "1h", "1d"],
            "1mo":  ["15m", "30m", "1h", "1d", "1wk"],
            "3mo":  ["1h", "1d", "1wk"],
            "6mo":  ["1d", "1wk"],
            "1y":   ["1d", "1wk", "1mo"],
            "2y":   ["1d", "1wk", "1mo"],
            "5y":   ["1d", "1wk", "1mo"],
            "10y":  ["1wk", "1mo"],
            "ytd":  ["1d", "1wk", "1mo"],
            "max":  ["1wk", "1mo"],
        }
        available_codes  = interval_map.get(period, ["1d"])
        available_labels = [INTERVAL_LABELS[c] for c in available_codes]
        interval_label_to_code = {INTERVAL_LABELS[c]: c for c in available_codes}
        interval_label = st.selectbox(
            "Intervalle", available_labels,
            help="Granularité des bougies. Les intervalles courts ne sont disponibles que sur des périodes courtes.",
        )
        interval = interval_label_to_code[interval_label]

        st.markdown("---")
        st.markdown("### Overlays graphique")
        show_ma20   = st.checkbox(
            "MA20", value=False,
            help="Moyenne mobile 20 jours — reflète la tendance court terme. Utilisée pour le swing trading et comme support/résistance dynamique.",
        )
        show_ma50   = st.checkbox(
            "MA50", value=True,
            help="Moyenne mobile 50 jours — tendance moyen terme. Référence clé des traders institutionnels. Un croisement MA50/MA200 forme un Golden Cross (haussier) ou Death Cross (baissier).",
        )
        show_ma100  = st.checkbox(
            "MA100", value=False,
            help="Moyenne mobile 100 jours — filtre intermédiaire entre le moyen et le long terme. Moins standard que la MA50 ou MA200.",
        )
        show_ma200  = st.checkbox(
            "MA200", value=False,
            help="Moyenne mobile 200 jours — tendance long terme. Un actif au-dessus de sa MA200 est généralement considéré en tendance haussière structurelle.",
        )
        show_bb     = st.checkbox(
            "Bollinger Bands", value=False,
            help="Bandes de Bollinger (MA20 ± 2σ) — mesurent la volatilité. Un prix touchant la bande haute/basse peut signaler un excès ou une cassure de volatilité.",
        )
        show_volume = st.checkbox(
            "Volume", value=True,
            help="Nombre de titres échangés par bougie. Un mouvement de prix confirmé par un volume élevé est plus fiable.",
        )
        log_scale   = st.checkbox(
            "Échelle log", value=False,
            help="Échelle logarithmique — utile sur longue période pour visualiser les rendements en % plutôt qu'en valeur absolue.",
        )

        chart_options = {
            "MA20": show_ma20, "MA50": show_ma50,
            "MA100": show_ma100, "MA200": show_ma200,
            "bollinger": show_bb, "show_volume": show_volume,
            "log_scale": log_scale,
            "interval": interval,
        }

        st.markdown("---")
        st.markdown("### Comparaison multi-actifs")
        compare_input = st.text_input(
            "Ajouter tickers (virgule)",
            placeholder="MSFT, GOOGL, AMZN",
        )

        st.markdown("---")
        st.markdown("### Sections à afficher")
        show_overview   = st.checkbox("Overview",             value=True)
        show_fund       = st.checkbox("Données fondamentales",value=True)
        show_charts     = st.checkbox("Graphiques avancés",   value=True)
        show_tech       = st.checkbox("Analyse technique",    value=True)
        show_stats      = st.checkbox("Statistiques",         value=True)
        show_risk       = st.checkbox("Analyse de risque",    value=True)
        show_compare    = st.checkbox("Comparaison",          value=True)
        show_export     = st.checkbox("Export données",       value=True)

    # ─── Data loading ────────────────────────────────────────
    if not symbol_input:
        st.info("Entrez un ticker dans la barre latérale pour commencer.")
        return

    with st.spinner(f"Chargement des données pour **{symbol_input}**…"):
        df_raw = load_data(symbol_input, period, interval)

    if df_raw.empty:
        st.error(f"Impossible de récupérer les données pour `{symbol_input}`. "
                 "Vérifiez le ticker et réessayez.")
        return

    df = compute_indicators(df_raw.copy())
    info   = get_ticker_info(symbol_input)
    stats  = compute_statistics(df)
    risk   = compute_risk_metrics(df)
    # Fear & Greed : indicateur CNN officiel (indépendant du ticker)
    fg     = fetch_fear_greed_cnn()

    # ─── Company name ────────────────────────────────────────
    company_name = info.get("longName") or info.get("shortName") or symbol_input
    sector       = info.get("sector", "")
    industry     = info.get("industry", "")
    currency     = info.get("currency", "USD")

    # ══════════════════════════════════════════════════════════
    #  TAB 1 — OVERVIEW
    # ══════════════════════════════════════════════════════════
    tabs = st.tabs([
        "Overview",
        "Fondamentaux",
        "Graphiques",
        "Technique",
        "Statistiques",
        "Risque",
        "Comparaison",
        "Export PPT",
        "Export Data",
        "Recherche Tickers",
        "Glossaire",
    ])

    # ── Tab: Overview ────────────────────────────────────────
    with tabs[0]:
        if not show_overview:
            st.info("Section désactivée dans la barre latérale.")
        else:
            # KPI row
            last_close  = df["Close"].iloc[-1]
            prev_close  = df["Close"].iloc[-2] if len(df) > 1 else last_close
            day_chg     = last_close - prev_close
            day_chg_pct = day_chg / prev_close * 100 if prev_close else 0
            total_ret   = stats.get("total_return", 0)
            vol_ann     = stats.get("volatility_annual", 0)
            max_dd      = stats.get("max_drawdown", 0)

            st.markdown(f"### {symbol_input}  <span style='font-size:0.75em;color:#8b92a5;font-weight:400;'>{company_name}</span>",
                        unsafe_allow_html=True)
            if sector:
                st.caption(f"{sector} · {industry}")

            col1, col2, col3, col4, col5 = st.columns([1.3, 1, 1, 1, 0.9])
            with col1:
                st.metric("Prix actuel", f"{last_close:.2f} {currency}",
                          f"{day_chg:+.2f} ({day_chg_pct:+.2f}%)")
            with col2:
                st.metric("Perf. période", fmt(total_ret*100, 2, "%"))
            with col3:
                st.metric("Vol. annualisée", fmt(vol_ann*100, 2, "%"))
            with col4:
                st.metric("Max Drawdown", fmt(max_dd*100, 2, "%"))
            with col5:
                render_fear_greed(fg)

            st.markdown("---")
            # Mini price chart
            fig_mini = create_price_chart(df, symbol_input, chart_options)
            st.plotly_chart(fig_mini, use_container_width=True, key="chart_1")

            # Company description
            desc = info.get("longBusinessSummary")
            if desc:
                with st.expander("Description de la société"):
                    # Découper en paragraphes pour aérer le texte
                    sentences = desc.replace(". ", ".\n\n").split("\n\n")
                    # Regrouper par blocs de ~2-3 phrases
                    paragraphs = []
                    buf = []
                    for s in sentences:
                        buf.append(s.strip())
                        if len(buf) >= 3:
                            paragraphs.append(" ".join(buf))
                            buf = []
                    if buf:
                        paragraphs.append(" ".join(buf))
                    for p in paragraphs:
                        st.markdown(f"<p style='line-height:1.7;color:#c5cae9;font-size:0.9rem;'>{p}</p>",
                                    unsafe_allow_html=True)

    # ── Tab: Fondamentaux ────────────────────────────────────
    with tabs[1]:
        if not show_fund:
            st.info("Section désactivée.")
        else:
            display_fundamentals(info)

            fund_data = get_fundamentals(symbol_input)

            # ── Recommandations analystes ─────────────────────
            reco = fund_data.get("recommendations")
            if reco is not None and not reco.empty:
                st.markdown('<p class="section-title">Recommandations Analystes</p>',
                            unsafe_allow_html=True)
                st.dataframe(reco.tail(20), use_container_width=True)

            # ── Analyst price targets ─────────────────────────
            apt = fund_data.get("analyst_price_targets")
            if apt is not None:
                st.markdown('<p class="section-title">Objectifs de Prix Analystes</p>',
                            unsafe_allow_html=True)
                if isinstance(apt, dict):
                    cols = st.columns(4)
                    for i, (k, v) in enumerate(apt.items()):
                        with cols[i % 4]:
                            st.metric(k, fmt(v, 2, "$") if isinstance(v, (int, float)) else str(v))
                elif isinstance(apt, pd.DataFrame):
                    st.dataframe(apt, use_container_width=True)

            # ── Dividendes historiques ────────────────────────
            divs = fund_data.get("dividends")
            if divs is not None and not divs.empty:
                st.markdown('<p class="section-title">Historique Dividendes</p>',
                            unsafe_allow_html=True)
                div_df = divs.reset_index()
                div_df.columns = ["Date", "Dividende"]
                div_df["Date"] = pd.to_datetime(div_df["Date"]).dt.strftime("%Y-%m-%d")
                fig_div = go.Figure(go.Bar(
                    x=div_df["Date"], y=div_df["Dividende"],
                    marker_color=UP_COLOR,
                ))
                fig_div.update_layout(
                    title="Dividendes distribués",
                    paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
                    xaxis=dict(gridcolor=GRID_COLOR, tickfont_color="#8b92a5"),
                    yaxis=dict(gridcolor=GRID_COLOR, tickfont_color="#8b92a5"),
                    height=300, margin=dict(l=10, r=10, t=40, b=10),
                )
                st.plotly_chart(fig_div, use_container_width=True, key="chart_2")

            # ── Institutional holders ─────────────────────────
            inst = fund_data.get("institutional_holders")
            if inst is not None and not inst.empty:
                st.markdown('<p class="section-title">Actionnaires Institutionnels</p>',
                            unsafe_allow_html=True)
                st.dataframe(inst.head(15), use_container_width=True)

            # ── Sustainability ────────────────────────────────
            sus = fund_data.get("sustainability")
            if sus is not None and not sus.empty:
                st.markdown('<p class="section-title">Score ESG</p>',
                            unsafe_allow_html=True)
                st.dataframe(sus, use_container_width=True)

            # ── Financials ────────────────────────────────────
            fin = fund_data.get("financials")
            if fin is not None and not fin.empty:
                st.markdown('<p class="section-title">États Financiers Annuels</p>',
                            unsafe_allow_html=True)
                st.dataframe(fin.style.format("{:,.0f}"), use_container_width=True)

    # ── Tab: Graphiques ──────────────────────────────────────
    with tabs[2]:
        if not show_charts:
            st.info("Section désactivée.")
        else:
            st.markdown('<p class="section-title">Graphique Principal</p>',
                        unsafe_allow_html=True)
            fig_main = create_price_chart(df, symbol_input, chart_options)
            st.plotly_chart(fig_main, use_container_width=True, key="chart_3")

            st.markdown('<p class="section-title">Drawdown</p>',
                        unsafe_allow_html=True)
            st.plotly_chart(create_drawdown_chart(df, symbol_input, interval),
                            use_container_width=True, key="chart_drawdown")

            st.markdown('<p class="section-title">Distribution des rendements</p>',
                        unsafe_allow_html=True)
            st.plotly_chart(create_returns_distribution(df),
                            use_container_width=True, key="chart_x4")

            st.markdown('<p class="section-title">Heatmap Rendements Mensuels</p>',
                        unsafe_allow_html=True)
            st.plotly_chart(create_monthly_returns_heatmap(df, symbol_input),
                            use_container_width=True, key="chart_x5")

    # ── Tab: Technique ───────────────────────────────────────
    with tabs[3]:
        if not show_tech:
            st.info("Section désactivée.")
        else:
            display_technical_analysis(df)
            st.markdown("---")
            st.markdown('<p class="section-title">Indicateurs Techniques</p>',
                        unsafe_allow_html=True)
            st.plotly_chart(create_technical_charts(df, interval), use_container_width=True, key="chart_4")

            # Dernières valeurs
            st.markdown('<p class="section-title">Dernières Valeurs</p>',
                        unsafe_allow_html=True)
            indicator_cols = [c for c in [
                "Close", "RSI", "MACD", "MACD_Signal",
                "BB_Upper", "BB_Mid", "BB_Lower", "ATR",
                "Stoch_K", "Stoch_D", "SMA20", "SMA50",
                "SMA200", "OBV",
            ] if c in df.columns]
            st.dataframe(
                df[indicator_cols].tail(20).style.format("{:,.2f}"),
                use_container_width=True,
            )

    # ── Tab: Statistiques ────────────────────────────────────
    with tabs[4]:
        if not show_stats:
            st.info("Section désactivée.")
        else:
            display_statistics(stats)
            st.markdown("---")
            fig_cum = create_cumulative_returns_chart({symbol_input: df})
            st.plotly_chart(fig_cum, use_container_width=True, key="chart_5")

            st.markdown('<p class="section-title">Rendements Périodiques</p>',
                        unsafe_allow_html=True)
            # Calculer rendements par période
            close = df["Close"]
            periods_ret = {}
            for label, n in [("1 sem.", 5), ("1 mois", 21), ("3 mois", 63),
                              ("6 mois", 126), ("1 an", 252), ("YTD", None)]:
                if label == "YTD":
                    start_year = close[close.index.year == close.index[-1].year].iloc[0]
                    periods_ret[label] = (close.iloc[-1] / start_year - 1) * 100
                elif len(close) > n:
                    periods_ret[label] = (close.iloc[-1] / close.iloc[-n] - 1) * 100
            if periods_ret:
                df_periods = pd.DataFrame.from_dict(
                    periods_ret, orient="index", columns=["Rendement (%)"]
                )
                fig_per = go.Figure(go.Bar(
                    x=df_periods.index,
                    y=df_periods["Rendement (%)"],
                    marker_color=[UP_COLOR if v >= 0 else DOWN_COLOR
                                  for v in df_periods["Rendement (%)"]],
                ))
                fig_per.update_layout(
                    title="Rendements par période (%)",
                    paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
                    yaxis=dict(ticksuffix="%", gridcolor=GRID_COLOR,
                               tickfont_color="#8b92a5"),
                    xaxis=dict(gridcolor=GRID_COLOR, tickfont_color="#8b92a5"),
                    height=320, margin=dict(l=10, r=10, t=40, b=10),
                )
                st.plotly_chart(fig_per, use_container_width=True, key="chart_6")

    # ── Tab: Risque ──────────────────────────────────────────
    with tabs[5]:
        if not show_risk:
            st.info("Section désactivée.")
        else:
            display_risk_analysis(risk)
            st.markdown("---")
            st.markdown('<p class="section-title">Rolling Volatility (30j)</p>',
                        unsafe_allow_html=True)
            if "Daily_Return" in df.columns:
                roll_vol = df["Daily_Return"].rolling(30).std() * np.sqrt(252) * 100
                fig_vol = go.Figure(go.Scatter(
                    x=df.index, y=roll_vol,
                    fill="tozeroy",
                    line=dict(color="#ff9800", width=1.5),
                    fillcolor="rgba(255,152,0,0.15)",
                    name="Vol. glissante 30j",
                ))
                _n = len(df)
                _step = max(1, _n // 25)
                _tv = list(df.index[::_step])
                if interval in ("1m", "5m", "15m", "30m"):
                    _tt = [t.strftime("%d %b %Hh%M") for t in _tv]
                elif interval == "1h":
                    _tt = [t.strftime("%d %b %Hh") for t in _tv]
                else:
                    _tt = [t.strftime("%d %b %Y") for t in _tv]
                fig_vol.update_layout(
                    paper_bgcolor=CHART_BG, plot_bgcolor=CHART_BG,
                    yaxis=dict(ticksuffix="%", gridcolor=GRID_COLOR,
                               tickfont_color="#8b92a5"),
                    xaxis=dict(type="category", tickvals=_tv, ticktext=_tt,
                               tickangle=-45, gridcolor=GRID_COLOR,
                               tickfont_color="#8b92a5"),
                    height=320, margin=dict(l=10, r=10, t=10, b=10),
                )
                st.plotly_chart(fig_vol, use_container_width=True, key="chart_7")

    # ── Tab: Comparaison ─────────────────────────────────────
    with tabs[6]:
        if not show_compare:
            st.info("Section désactivée.")
        else:
            tickers_compare = [symbol_input]
            if compare_input:
                extras = [t.strip().upper() for t in compare_input.split(",") if t.strip()]
                tickers_compare.extend(extras[:5])

            if len(tickers_compare) == 1:
                st.info("Ajoutez d'autres tickers dans la barre latérale pour comparer.")
            else:
                dfs_compare = {symbol_input: df}
                with st.spinner("Chargement des actifs comparables…"):
                    for sym in tickers_compare[1:]:
                        d = load_data(sym, period, interval)
                        if not d.empty:
                            dfs_compare[sym] = compute_indicators(d.copy())

                st.plotly_chart(
                    create_cumulative_returns_chart(dfs_compare),
                    use_container_width=True,
                    key="chart_compare_cum",
                )
                if len(dfs_compare) >= 2:
                    st.plotly_chart(
                        create_correlation_heatmap(dfs_compare),
                        use_container_width=True,
                        key="chart_compare_corr",
                    )

                # Tableau comparatif
                st.markdown('<p class="section-title">Métriques Comparées</p>',
                            unsafe_allow_html=True)
                rows_cmp = []
                for sym, d in dfs_compare.items():
                    s = compute_statistics(d)
                    r = compute_risk_metrics(d)
                    i = get_ticker_info(sym)
                    rows_cmp.append({
                        "Ticker":           sym,
                        "Rend. Total (%)":  fmt(s.get("total_return", 0)*100, 2),
                        "Vol. Ann. (%)":    fmt(s.get("volatility_annual", 0)*100, 2),
                        "Sharpe":           fmt(s.get("sharpe_ratio", 0), 2),
                        "Max DD (%)":       fmt(s.get("max_drawdown", 0)*100, 2),
                        "VaR 95% (%)":      fmt(r.get("VaR_95", 0)*100, 2),
                        "Beta":             fmt(i.get("beta"), 2),
                        "Market Cap":       fmt_large(i.get("marketCap"), i.get("currency", "")),
                    })
                st.dataframe(pd.DataFrame(rows_cmp).set_index("Ticker"),
                             use_container_width=True)

    # ── Tab: Export ──────────────────────────────────────────
    with tabs[7]:
        st.markdown('<p class="section-title">Export PowerPoint</p>', unsafe_allow_html=True)
        st.markdown(
            f"Générez un support de présentation complet pour **{company_name}** (`{symbol_input}`) "
            f"sur la période **{period}** ({df.index[0].strftime('%d/%m/%Y')} — {df.index[-1].strftime('%d/%m/%Y')})."
        )
        st.markdown("""
        Le fichier PowerPoint contiendra :
        - Slide de titre avec les informations du titre
        - KPIs clés : prix, performance, volatilité, Fear & Greed…
        - Graphique de prix & volume avec overlays
        - Indicateurs techniques (RSI, MACD, Stochastique, OBV)
        - Statistiques de performance (Sharpe, Sortino, Skewness, Kurtosis…)
        - Distribution des rendements vs loi normale
        - Drawdown historique
        - Heatmap des rendements mensuels
        - Métriques de risque (VaR, CVaR…)
        - Slide de disclaimer
        """)

        if st.button("Générer le PowerPoint", type="primary", key="btn_pptx"):
            with st.spinner("Génération du fichier PowerPoint en cours…"):
                try:
                    pptx_bytes = generate_pptx(
                        symbol=symbol_input,
                        company_name=company_name,
                        sector=sector,
                        industry=industry,
                        currency=currency,
                        period=period,
                        info=info,
                        df=df,
                        stats=stats,
                        risk=risk,
                        fg=fg,
                        chart_options=chart_options,
                    )
                    st.success("PowerPoint généré avec succès.")
                    st.download_button(
                        label="Télécharger le PowerPoint",
                        data=pptx_bytes,
                        file_name=f"{symbol_input}_Apex_Markets_{period}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="dl_pptx",
                    )
                except Exception as e:
                    st.error(f"Erreur lors de la génération : {e}")

    # ── Tab: Export Data ─────────────────────────────────────
    with tabs[8]:
        if not show_export:
            st.info("Section désactivée.")
        else:
            st.markdown('<p class="section-title">Export des Données</p>',
                        unsafe_allow_html=True)
            st.markdown("Téléchargez les données brutes et indicateurs calculés.")

            col_a, col_b = st.columns(2)
            with col_a:
                csv_raw = df_raw.to_csv().encode("utf-8")
                st.download_button(
                    "⬇️ Données OHLCV (CSV)",
                    csv_raw,
                    file_name=f"{symbol_input}_ohlcv_{period}.csv",
                    mime="text/csv",
                )
            with col_b:
                csv_ind = df.to_csv().encode("utf-8")
                st.download_button(
                    "⬇️ Données + Indicateurs (CSV)",
                    csv_ind,
                    file_name=f"{symbol_input}_indicators_{period}.csv",
                    mime="text/csv",
                )

            # Stats JSON export
            stats_export = {k: (v if not isinstance(v, float) or not np.isnan(v) else None)
                            for k, v in stats.items()}
            risk_export  = {k: (v if not isinstance(v, float) or not np.isnan(v) else None)
                            for k, v in risk.items()}
            export_json = json.dumps(
                {"symbol": symbol_input, "period": period,
                 "statistics": stats_export, "risk": risk_export},
                indent=2,
            )
            st.download_button(
                "⬇️ Statistiques & Risque (JSON)",
                export_json,
                file_name=f"{symbol_input}_stats_{period}.json",
                mime="application/json",
            )

            # Preview
            st.markdown('<p class="section-title">Aperçu des données</p>',
                        unsafe_allow_html=True)
            st.dataframe(df.tail(30), use_container_width=True)

    # ── Tab: Recherche Tickers ───────────────────────────────
    with tabs[9]:
        st.title("Recherche de Tickers")
        st.caption("Tape un nom d'entreprise, un indice, une crypto ou même une approximation.")

        recherche = st.text_input(
            "Rechercher",
            placeholder="Ex : Total, CAC 40, bitcoin, LVMH, S&P 500…",
        ).strip()

        if recherche:
            with st.spinner("Recherche en cours…"):
                resultats = rechercher_tickers(recherche)

            if not resultats:
                st.warning("Aucun résultat. Essaie une autre orthographe.")
            else:
                st.success(f"{len(resultats)} résultat(s) trouvé(s)")
                st.divider()
                for r in resultats:
                    c1, c2, c3 = st.columns([1.5, 4, 1.5])
                    c1.markdown(f"### `{r['symbol']}`")
                    c2.markdown(f"**{r['name']}**")
                    c2.caption(
                        f"{TYPE_LABELS.get(r['type'], r['type'])} · {r['exchange']}"
                    )
                    if c3.button("Copier", key=f"copy_{r['symbol']}"):
                        st.toast(f"Ticker **{r['symbol']}** — copiez-le dans le champ Ticker de la sidebar.")
                    st.divider()
        else:
            st.markdown("""
            ### Exemples de recherches
            | Ce que tu tapes | Ticker trouvé |
            |---|---|
            | `Apple` | `AAPL` |
            | `Total` | `TTE.PA` |
            | `LVMH` | `MC.PA` |
            | `BNP` | `BNP.PA` |
            | `CAC` | `^FCHI` |
            | `S&P 500` | `^GSPC` |
            | `bitcoin` | `BTC-USD` |
            | `Tesla` | `TSLA` |
            | `Nasdaq` | `^IXIC` |
            | `Gold` | `GC=F` |
            | `Euro Dollar` | `EURUSD=X` |

            > **Astuce :** Une fois le ticker trouvé, copiez-le et collez-le dans le champ **Ticker principal** de la barre latérale.
            """)

    # ── Tab: Glossaire ───────────────────────────────────────
    with tabs[10]:
        st.title("Glossaire Financier")
        st.caption("Définitions simples de toutes les métriques affichées dans l'application.")

        # Filtre de recherche dans le glossaire
        filtre = st.text_input("Filtrer le glossaire", placeholder="Ex : P/E, Sharpe, Beta…").strip().lower()

        categories = {
            "Prix & Marché": [
                "Prix actuel", "Open", "High (jour)", "Low (jour)",
                "Clôture préc.", "52-sem. High", "52-sem. Low", "Volume",
            ],
            "Valorisation": [
                "Market Cap", "Enterprise V.", "Beta", "Float",
                "P/E (TTM)", "P/E Forward", "PEG", "P/B", "P/S",
                "EV/EBITDA", "EV/Rev.", "52w perf",
            ],
            "Rentabilité": [
                "ROE", "ROA", "Marge brute", "Marge opérat.", "Marge nette",
                "EBITDA", "Chiffre d'aff.", "Rev. Growth",
            ],
            "Dividendes": [
                "Dividend Yield", "Dividende/Act", "Payout Ratio", "Ex-Date",
            ],
            "Structure Financière": [
                "Dette/Cap. prop.", "Total Cash", "Total Dette",
                "Free Cash Flow", "Op. Cash Flow", "EPS (TTM)", "EPS Forward", "Earn. Growth",
            ],
            "Performance & Statistiques": [
                "Rend. total", "Rend. moy. annuel", "Vol. annualisée",
                "Sharpe Ratio", "Sortino Ratio", "Max Drawdown",
                "Jours positifs", "Asymétrie", "Kurtosis",
                "Meilleure jour", "Pire jour",
            ],
            "Risque": [
                "VaR 95% (1j)", "VaR 99% (1j)", "CVaR 95%", "CVaR 99%", "Downside Risk",
            ],
        }

        for categorie, metriques in categories.items():
            # Filtrer selon la recherche
            metriques_filtrees = [
                m for m in metriques
                if not filtre or filtre in m.lower() or filtre in TOOLTIPS.get(m, "").lower()
            ]
            if not metriques_filtrees:
                continue

            st.markdown(f'<p class="section-title">{categorie}</p>', unsafe_allow_html=True)
            for metrique in metriques_filtrees:
                definition = TOOLTIPS.get(metrique, "Définition non disponible.")
                col_a, col_b = st.columns([1, 3])
                with col_a:
                    st.markdown(f"**{metrique}**")
                with col_b:
                    st.markdown(f"<span style='color:#c5cae9;'>{definition}</span>",
                                unsafe_allow_html=True)
            st.markdown("---")

        st.markdown("""
        > **Rappel important** : Ces métriques sont des outils d'aide à la décision,
        > pas des garanties. Tout investissement comporte un risque de perte en capital.
        """)

    # ─── Footer ──────────────────────────────────────────────
    st.markdown("---")
    st.caption(
        "*Apex Markets est un projet personnel à vocation pédagogique et ne constitue en aucun cas "
        "un conseil en investissement, une recommandation d'achat ou de vente, ni une sollicitation "
        "d'opérations sur instruments financiers. Les données de marché sont fournies par Yahoo Finance "
        "via la bibliothèque open-source yfinance ; elles peuvent être différées, incomplètes ou inexactes. "
        "L'auteur décline toute responsabilité quant aux décisions prises sur la base des informations "
        "affichées. Tout investissement comporte un risque de perte en capital.*"
    )


# ══════════════════════════════════════════════════════════════
#  SECTION 7 — PPT EXPORT
# ══════════════════════════════════════════════════════════════

def fig_to_base64(fig: go.Figure, width: int = 900, height: int = 420) -> str:
    """Convertit un figure Plotly en base64 PNG via kaleido (fallback: None)."""
    try:
        import base64 as _b64
        img_bytes = fig.to_image(format="png", width=width, height=height, scale=2)
        return "image/png;base64," + _b64.b64encode(img_bytes).decode()
    except Exception:
        return None

# ── Graphiques matplotlib pour export PPT (sans kaleido) ─────────────────

def _mpl_buf_to_b64(buf) -> str:
    import base64
    buf.seek(0)
    return "image/png;base64," + base64.b64encode(buf.read()).decode()


def mpl_price_chart(df: pd.DataFrame, symbol: str) -> str:
    """Cours de clôture + volume en matplotlib."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.gridspec as gridspec
    import io

    fig = plt.figure(figsize=(13, 6), facecolor="#0e1117")
    gs  = gridspec.GridSpec(2, 1, height_ratios=[3, 1], hspace=0.05)
    ax1 = fig.add_subplot(gs[0])
    ax2 = fig.add_subplot(gs[1], sharex=ax1)

    close = df["Close"]
    ax1.plot(df.index, close, color="#00aaff", linewidth=1.4, label="Clôture")

    # MAs
    for col, color in [("SMA20","#f9a825"),("SMA50","#42a5f5"),("SMA200","#ef5350")]:
        if col in df.columns:
            ax1.plot(df.index, df[col], color=color, linewidth=1, alpha=0.8, label=col)

    # Bollinger
    if "BB_Upper" in df.columns:
        ax1.fill_between(df.index, df["BB_Lower"], df["BB_Upper"],
                         alpha=0.08, color="#ff9800")
        ax1.plot(df.index, df["BB_Upper"], color="#ff9800", linewidth=0.7, linestyle="--")
        ax1.plot(df.index, df["BB_Lower"], color="#ff9800", linewidth=0.7, linestyle="--")

    # Volume
    colors = ["#00e676" if df["Close"].iloc[i] >= df["Open"].iloc[i] else "#ff5252"
              for i in range(len(df))]
    ax2.bar(df.index, df["Volume"], color=colors, alpha=0.7, width=0.8)

    for ax in [ax1, ax2]:
        ax.set_facecolor("#0e1117")
        ax.tick_params(colors="#8b92a5", labelsize=7)
        ax.spines[:].set_color("#2a2d3e")
        ax.grid(color="#2a2d3e", linewidth=0.5)
        ax.yaxis.label.set_color("#8b92a5")

    ax1.set_title(f"{symbol} — Prix & Volume", color="#00d4ff", fontsize=12, pad=8)
    ax1.legend(fontsize=7, facecolor="#1e2130", labelcolor="#ccc", framealpha=0.8)
    ax2.set_ylabel("Volume", color="#8b92a5", fontsize=7)
    plt.setp(ax1.get_xticklabels(), visible=False)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="#0e1117")
    plt.close(fig)
    return _mpl_buf_to_b64(buf)


def mpl_indicators_chart(df: pd.DataFrame) -> str:
    """RSI + MACD + Stochastique en matplotlib."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.gridspec as gridspec
    import io

    fig = plt.figure(figsize=(13, 8), facecolor="#0e1117")
    gs  = gridspec.GridSpec(3, 1, hspace=0.35)
    axes = [fig.add_subplot(gs[i]) for i in range(3)]

    # RSI
    if "RSI" in df.columns:
        axes[0].plot(df.index, df["RSI"], color="#00aaff", linewidth=1.3)
        axes[0].axhline(70, color="#ff5252", linestyle="--", linewidth=0.8, alpha=0.7)
        axes[0].axhline(30, color="#00e676", linestyle="--", linewidth=0.8, alpha=0.7)
        axes[0].fill_between(df.index, 30, 70, alpha=0.04, color="white")
        axes[0].set_ylim(0, 100)
        axes[0].set_title("RSI (14)", color="#00d4ff", fontsize=9)

    # MACD
    if "MACD" in df.columns:
        axes[1].plot(df.index, df["MACD"],        color="#42a5f5", linewidth=1.2, label="MACD")
        axes[1].plot(df.index, df["MACD_Signal"], color="#f9a825", linewidth=1.2, label="Signal")
        hist = df["MACD_Hist"].fillna(0)
        colors = ["#00e676" if v >= 0 else "#ff5252" for v in hist]
        axes[1].bar(df.index, hist, color=colors, alpha=0.6, width=0.8)
        axes[1].set_title("MACD (12/26/9)", color="#00d4ff", fontsize=9)
        axes[1].legend(fontsize=7, facecolor="#1e2130", labelcolor="#ccc", framealpha=0.8)

    # Stochastique
    if "Stoch_K" in df.columns:
        axes[2].plot(df.index, df["Stoch_K"], color="#ce93d8", linewidth=1.2, label="%K")
        axes[2].plot(df.index, df["Stoch_D"], color="#f9a825", linewidth=1, linestyle="--", label="%D")
        axes[2].axhline(80, color="#ff5252", linestyle="--", linewidth=0.8, alpha=0.7)
        axes[2].axhline(20, color="#00e676", linestyle="--", linewidth=0.8, alpha=0.7)
        axes[2].set_ylim(0, 100)
        axes[2].set_title("Stochastique (14,3)", color="#00d4ff", fontsize=9)
        axes[2].legend(fontsize=7, facecolor="#1e2130", labelcolor="#ccc", framealpha=0.8)

    for ax in axes:
        ax.set_facecolor("#0e1117")
        ax.tick_params(colors="#8b92a5", labelsize=7)
        ax.spines[:].set_color("#2a2d3e")
        ax.grid(color="#2a2d3e", linewidth=0.5)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="#0e1117")
    plt.close(fig)
    return _mpl_buf_to_b64(buf)


def mpl_distribution_chart(df: pd.DataFrame) -> str:
    """Distribution des rendements + loi normale en matplotlib."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import io

    fig, ax = plt.subplots(figsize=(10, 5), facecolor="#0e1117")
    ax.set_facecolor("#0e1117")

    ret = df["Daily_Return"].dropna() * 100
    mean_r, std_r = ret.mean(), ret.std()
    skew_r, kurt_r = ret.skew(), ret.kurt()

    ax.hist(ret, bins=60, density=True, color="#00aaff", alpha=0.75, label="Rendements")

    x = np.linspace(ret.min(), ret.max(), 300)
    from scipy import stats as sp
    ax.plot(x, sp.norm.pdf(x, mean_r, std_r), color="#ff9800", linewidth=2, label="Loi Normale")

    ax.axvline(mean_r, color="#f9a825", linestyle="--", linewidth=1.2, label=f"μ={mean_r:.2f}%")
    for k, alpha in [(1, 0.6), (2, 0.3)]:
        for sign in [-1, 1]:
            ax.axvline(mean_r + sign*k*std_r, color="#ab47bc", linestyle=":", linewidth=0.9, alpha=alpha)

    ax.set_title("Distribution des Rendements vs Loi Normale", color="#00d4ff", fontsize=11)
    ax.set_xlabel("Rendement (%)", color="#8b92a5", fontsize=9)
    ax.set_ylabel("Densité", color="#8b92a5", fontsize=9)
    ax.tick_params(colors="#8b92a5", labelsize=8)
    ax.spines[:].set_color("#2a2d3e")
    ax.grid(color="#2a2d3e", linewidth=0.5)
    ax.legend(fontsize=8, facecolor="#1e2130", labelcolor="#ccc", framealpha=0.8)

    ax.text(0.01, 0.97, f"Skewness: {skew_r:.3f}   Kurtosis: {kurt_r:.3f}",
            transform=ax.transAxes, color="#c5cae9", fontsize=8,
            verticalalignment="top",
            bbox=dict(boxstyle="round", facecolor="#1e2130", alpha=0.8))

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="#0e1117")
    plt.close(fig)
    return _mpl_buf_to_b64(buf)


def mpl_drawdown_chart(df: pd.DataFrame, symbol: str) -> str:
    """Drawdown historique en matplotlib."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import io

    fig, ax = plt.subplots(figsize=(13, 4), facecolor="#0e1117")
    ax.set_facecolor("#0e1117")

    if "Drawdown" in df.columns:
        dd = df["Drawdown"] * 100
        ax.fill_between(df.index, dd, 0, color="#ff5252", alpha=0.4)
        ax.plot(df.index, dd, color="#ff5252", linewidth=1.2)

    ax.set_title(f"{symbol} — Drawdown Historique (%)", color="#00d4ff", fontsize=11)
    ax.set_ylabel("%", color="#8b92a5", fontsize=9)
    ax.tick_params(colors="#8b92a5", labelsize=8)
    ax.spines[:].set_color("#2a2d3e")
    ax.grid(color="#2a2d3e", linewidth=0.5)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="#0e1117")
    plt.close(fig)
    return _mpl_buf_to_b64(buf)


def mpl_monthly_heatmap(df: pd.DataFrame, symbol: str) -> str:
    """Heatmap des rendements mensuels en matplotlib."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import io

    if df.empty or "Daily_Return" not in df.columns:
        return None

    monthly = df["Daily_Return"].resample("ME").apply(lambda x: (1+x).prod()-1) * 100
    monthly = monthly.to_frame("Return")
    monthly["Year"]  = monthly.index.year
    monthly["Month"] = monthly.index.month
    pivot = monthly.pivot_table(index="Year", columns="Month", values="Return")
    month_labels = ["Jan","Fév","Mar","Avr","Mai","Jun","Jul","Aoû","Sep","Oct","Nov","Déc"]
    pivot.columns = [month_labels[m-1] for m in pivot.columns]

    fig, ax = plt.subplots(figsize=(13, max(3, len(pivot)*0.5+1.5)), facecolor="#0e1117")
    ax.set_facecolor("#0e1117")

    import matplotlib.colors as mcolors
    cmap = plt.cm.RdYlGn
    vmax = max(abs(pivot.values[~np.isnan(pivot.values)]).max(), 1)
    norm = mcolors.TwoSlopeNorm(vmin=-vmax, vcenter=0, vmax=vmax)

    im = ax.imshow(pivot.values, cmap=cmap, norm=norm, aspect="auto")

    ax.set_xticks(range(len(pivot.columns)))
    ax.set_xticklabels(pivot.columns, color="#8b92a5", fontsize=8)
    ax.set_yticks(range(len(pivot.index)))
    ax.set_yticklabels(pivot.index.astype(str), color="#8b92a5", fontsize=8)
    ax.tick_params(colors="#8b92a5")
    ax.spines[:].set_color("#2a2d3e")

    for i in range(pivot.shape[0]):
        for j in range(pivot.shape[1]):
            val = pivot.values[i, j]
            if not np.isnan(val):
                ax.text(j, i, f"{val:.1f}%", ha="center", va="center",
                        fontsize=7, color="black" if abs(val) < vmax*0.6 else "white")

    plt.colorbar(im, ax=ax, label="%").ax.yaxis.label.set_color("#8b92a5")
    ax.set_title(f"{symbol} — Rendements Mensuels (%)", color="#00d4ff", fontsize=11)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="#0e1117")
    plt.close(fig)
    return _mpl_buf_to_b64(buf)


def generate_pptx(
    symbol: str,
    company_name: str,
    sector: str,
    industry: str,
    currency: str,
    period: str,
    info: dict,
    df: pd.DataFrame,
    stats: dict,
    risk: dict,
    fg: int,
    chart_options: dict,
    **kwargs,
) -> bytes:
    """
    Génère un fichier PPTX complet via python-pptx (pur Python, compatible Streamlit Cloud).
    Retourne les bytes du fichier .pptx.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import base64 as _b64

    # ── Couleurs ──
    BG      = RGBColor(0x0E, 0x11, 0x17)
    CARD    = RGBColor(0x1E, 0x21, 0x30)
    ACC     = RGBColor(0x00, 0x66, 0xFF)
    ACC2    = RGBColor(0x00, 0xD4, 0xFF)
    UP      = RGBColor(0x00, 0xE6, 0x76)
    DOWN    = RGBColor(0xFF, 0x52, 0x52)
    ORA     = RGBColor(0xFF, 0x98, 0x00)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED   = RGBColor(0x8B, 0x92, 0xA5)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]  # blank layout

    def _add_text(slide, text, left, top, width, height,
                  font_size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align

    def _add_rect(slide, left, top, width, height, fill_color):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()

    def _set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_card(slide, x, y, w, h, label, value, accent_color):
        _add_rect(slide, x, y, w, h, CARD)
        _add_rect(slide, x, y, 0.05, h, accent_color)
        _add_text(slide, label, x+0.12, y+0.05, w-0.2, 0.22, font_size=8, color=MUTED)
        _add_text(slide, value, x+0.12, y+0.28, w-0.2, 0.38, font_size=13, color=WHITE, bold=True)

    def _footer(slide):
        dr = f"{df.index[0].strftime('%d/%m/%Y')} — {df.index[-1].strftime('%d/%m/%Y')}"
        _add_text(slide, f"Apex Markets — {symbol} — {dr}",
                  0.3, 5.35, 9.4, 0.2, font_size=8, color=MUTED)

    def _add_chart_slide(title, img_b64):
        if not img_b64:
            return
        s = prs.slides.add_slide(blank)
        _set_bg(s, BG)
        _add_rect(s, 0, 0, 10, 0.55, CARD)
        _add_text(s, title, 0.3, 0.1, 9.4, 0.35, font_size=14, bold=True, color=ACC2)
        # Décoder l'image base64
        header, data = img_b64.split(",", 1)
        img_bytes = _b64.b64decode(data)
        img_stream = io.BytesIO(img_bytes)
        s.shapes.add_picture(img_stream, Inches(0.15), Inches(0.6), Inches(9.7), Inches(4.85))
        _footer(s)

    # ── Préparer les images ──
    charts = {}
    for key, fn in {
        "price":   lambda: mpl_price_chart(df, symbol),
        "tech":    lambda: mpl_indicators_chart(df),
        "dist":    lambda: mpl_distribution_chart(df),
        "dd":      lambda: mpl_drawdown_chart(df, symbol),
        "monthly": lambda: mpl_monthly_heatmap(df, symbol),
    }.items():
        try:
            charts[key] = fn()
        except Exception:
            charts[key] = None

    # ── Calculs ──
    last_close  = float(df["Close"].iloc[-1])
    prev_close  = float(df["Close"].iloc[-2]) if len(df) > 1 else last_close
    day_chg_pct = (last_close - prev_close) / prev_close * 100 if prev_close else 0.0
    total_ret   = stats.get("total_return", 0) * 100
    vol_ann     = stats.get("volatility_annual", 0) * 100
    sharpe      = stats.get("sharpe_ratio", 0)
    sortino     = stats.get("sortino_ratio", 0)
    max_dd      = stats.get("max_drawdown", 0) * 100
    skewness_v  = stats.get("skewness", 0)
    kurtosis_v  = stats.get("kurtosis", 0)
    pos_days    = stats.get("positive_days_pct", 0)
    var95       = risk.get("VaR_95", 0) * 100
    var99       = risk.get("VaR_99", 0) * 100
    cvar95      = risk.get("CVaR_95", 0) * 100
    downside    = risk.get("downside_dev", 0) * 100

    # ── SLIDE 1 : Titre ──
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG)
    _add_rect(s, 0, 0, 10, 0.08, ACC)
    _add_rect(s, 0, 5.545, 10, 0.08, ACC)
    subtitle = f"{sector}{' · ' + industry if industry else ''} · Période : {period}"
    _add_text(s, f"{symbol} — {company_name}", 0.5, 1.7, 9, 1.3,
              font_size=34, bold=True, color=ACC2, align=PP_ALIGN.CENTER)
    _add_text(s, subtitle, 0.5, 3.1, 9, 0.7,
              font_size=14, color=MUTED, align=PP_ALIGN.CENTER)
    _add_text(s, "Apex Markets", 0.5, 5.0, 9, 0.4,
              font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 : KPIs ──
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG)
    _add_rect(s, 0, 0, 10, 0.55, CARD)
    _add_text(s, f"Indicateurs Clés — {symbol}", 0.3, 0.1, 9.4, 0.35,
              font_size=14, bold=True, color=ACC2)
    row1 = [
        ("Prix actuel", f"{last_close:.2f} {currency}  ({day_chg_pct:+.2f}%)", UP if day_chg_pct >= 0 else DOWN),
        ("Perf. période", f"{total_ret:+.2f}%", UP if total_ret >= 0 else DOWN),
        ("Vol. annualisée", f"{vol_ann:.2f}%", ORA),
        ("Max Drawdown", f"{max_dd:.2f}%", DOWN),
    ]
    for i, (lbl, val, col) in enumerate(row1):
        _add_card(s, 0.2 + i*2.4, 0.65, 2.2, 0.85, lbl, val, col)
    row2 = [
        ("Sharpe Ratio", f"{sharpe:.3f}", ACC),
        ("Sortino Ratio", f"{sortino:.3f}", ACC),
        ("Market Cap", fmt_large(info.get("marketCap"), currency), ACC),
        ("Beta", fmt(info.get("beta"), 2), ACC2),
    ]
    for i, (lbl, val, col) in enumerate(row2):
        _add_card(s, 0.2 + i*2.4, 1.6, 2.2, 0.85, lbl, val, col)
    row3 = [
        ("P/E (TTM)", fmt(info.get("trailingPE"), 1), ACC2),
        ("Div. Yield", fmt(info.get("dividendYield"), 2, "%"), UP),
        ("ROE", fmt(info.get("returnOnEquity", 0)*100 if info.get("returnOnEquity") else None, 1, "%"), UP),
        ("Marge nette", fmt(info.get("profitMargins", 0)*100 if info.get("profitMargins") else None, 1, "%"), UP),
    ]
    for i, (lbl, val, col) in enumerate(row3):
        _add_card(s, 0.2 + i*2.4, 2.55, 2.2, 0.85, lbl, val, col)
    # Fear & Greed
    fg_color = DOWN if fg < 25 else (ORA if fg < 45 else (RGBColor(0xF9,0xA8,0x25) if fg < 55 else (UP if fg < 75 else UP)))
    fg_label = "Peur extrême" if fg < 25 else ("Peur" if fg < 45 else ("Neutre" if fg < 55 else ("Avidité" if fg < 75 else "Avidité extrême")))
    _add_rect(s, 0.2, 3.55, 9.6, 0.85, CARD)
    _add_text(s, f"Fear & Greed (CNN) : {fg}/100 — {fg_label}", 0.3, 3.6, 9, 0.35,
              font_size=12, bold=True, color=fg_color)
    _add_rect(s, 0.3, 4.0, 9.4, 0.15, RGBColor(0x2A, 0x2D, 0x3E))
    _add_rect(s, 0.3, 4.0, 9.4 * fg / 100, 0.15, fg_color)
    _footer(s)

    # ── SLIDE 3 : Profil société ──
    desc = info.get("longBusinessSummary", "")
    if desc or sector:
        s = prs.slides.add_slide(blank)
        _set_bg(s, BG)
        _add_rect(s, 0, 0, 10, 0.55, CARD)
        _add_text(s, f"Profil — {company_name}", 0.3, 0.1, 9.4, 0.35,
                  font_size=14, bold=True, color=ACC2)
        y_pos = 0.7
        profile_items = [
            ("Secteur", sector),
            ("Industrie", industry),
            ("Devise", currency),
            ("Place de cotation", info.get("exchange", "N/A")),
            ("Pays", info.get("country", "N/A")),
            ("Nb. employés", f"{info.get('fullTimeEmployees', 'N/A'):,}".replace(",", " ") if info.get("fullTimeEmployees") else "N/A"),
            ("Site web", info.get("website", "N/A")),
        ]
        for i, (lbl, val) in enumerate(profile_items):
            if val and val != "N/A":
                _add_text(s, f"{lbl} :  {val}", 0.4, y_pos + i * 0.32, 9.2, 0.28,
                          font_size=11, color=WHITE)
        if desc:
            # Tronquer à ~500 caractères pour tenir dans la slide
            short_desc = desc[:500] + ("…" if len(desc) > 500 else "")
            _add_text(s, short_desc, 0.4, 3.0, 9.2, 2.2, font_size=10, color=MUTED)
        _footer(s)

    # ── SLIDE 4 : Fondamentaux détaillés ──
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG)
    _add_rect(s, 0, 0, 10, 0.55, CARD)
    _add_text(s, f"Données Fondamentales — {symbol}", 0.3, 0.1, 9.4, 0.35,
              font_size=14, bold=True, color=ACC2)
    cur_sym = {"USD": "$", "EUR": "€", "GBP": "£", "JPY": "¥", "CHF": "CHF "}.get(currency, currency + " ")
    fund_items = [
        ("P/E (TTM)", fmt(info.get("trailingPE"), 1)),
        ("P/E Forward", fmt(info.get("forwardPE"), 1)),
        ("PEG", fmt(info.get("pegRatio"), 2)),
        ("P/B", fmt(info.get("priceToBook"), 2)),
        ("P/S", fmt(info.get("priceToSalesTrailing12Months"), 2)),
        ("EV/EBITDA", fmt(info.get("enterpriseToEbitda"), 1)),
        ("EV/Revenue", fmt(info.get("enterpriseToRevenue"), 2)),
        ("ROE", fmt(info.get("returnOnEquity", 0)*100 if info.get("returnOnEquity") else None, 1, "%")),
        ("ROA", fmt(info.get("returnOnAssets", 0)*100 if info.get("returnOnAssets") else None, 1, "%")),
        ("Marge brute", fmt(info.get("grossMargins", 0)*100 if info.get("grossMargins") else None, 1, "%")),
        ("Marge opérat.", fmt(info.get("operatingMargins", 0)*100 if info.get("operatingMargins") else None, 1, "%")),
        ("Marge nette", fmt(info.get("profitMargins", 0)*100 if info.get("profitMargins") else None, 1, "%")),
        ("Chiffre d'aff.", fmt_large(info.get("totalRevenue"), currency)),
        ("EBITDA", fmt_large(info.get("ebitda"), currency)),
        ("Free Cash Flow", fmt_large(info.get("freeCashflow"), currency)),
        ("Total Cash", fmt_large(info.get("totalCash"), currency)),
        ("Total Dette", fmt_large(info.get("totalDebt"), currency)),
        ("Dette/Cap. propres", fmt(info.get("debtToEquity"), 2)),
        ("Div. Yield", fmt(info.get("dividendYield"), 2, "%")),
        ("Payout Ratio", fmt(info.get("payoutRatio", 0)*100 if info.get("payoutRatio") else None, 1, "%")),
    ]
    for i, (lbl, val) in enumerate(fund_items):
        row = i // 4
        col = i % 4
        _add_card(s, 0.2 + col * 2.4, 0.65 + row * 0.95, 2.2, 0.8, lbl, val, ACC)
    _footer(s)

    # ── SLIDES 5-6 : Graphiques ──
    _add_chart_slide("Graphique de Prix & Volume", charts.get("price"))
    _add_chart_slide("Indicateurs Techniques — RSI · MACD · Stochastique · OBV", charts.get("tech"))

    # ── SLIDE 7 : Stats ──
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG)
    _add_rect(s, 0, 0, 10, 0.55, CARD)
    _add_text(s, "Performance & Statistiques", 0.3, 0.1, 9.4, 0.35,
              font_size=14, bold=True, color=ACC2)
    perf_data = [
        ("Rend. total", f"{total_ret:+.2f}%"), ("Vol. annualisée", f"{vol_ann:.2f}%"),
        ("Sharpe Ratio", f"{sharpe:.3f}"), ("Sortino Ratio", f"{sortino:.3f}"),
        ("Max Drawdown", f"{max_dd:.2f}%"), ("Jours positifs", f"{pos_days:.1f}%"),
        ("Skewness", f"{skewness_v:.3f}"), ("Kurtosis", f"{kurtosis_v:.3f}"),
    ]
    for i, (lbl, val) in enumerate(perf_data):
        _add_card(s, 0.2 + (i % 4) * 2.4, 0.65 + (i // 4) * 1.05, 2.2, 0.9, lbl, val, ACC)
    _footer(s)

    # ── SLIDES 6-8 : Graphiques ──
    _add_chart_slide("Distribution des Rendements — vs Loi Normale", charts.get("dist"))
    _add_chart_slide("Drawdown Historique", charts.get("dd"))
    _add_chart_slide("Rendements Mensuels", charts.get("monthly"))

    # ── SLIDE 9 : Risque ──
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG)
    _add_rect(s, 0, 0, 10, 0.55, CARD)
    _add_text(s, "Analyse de Risque", 0.3, 0.1, 9.4, 0.35,
              font_size=14, bold=True, color=ACC2)
    risk_data = [
        ("VaR 95% (1j)", f"{var95:.2f}%"), ("VaR 99% (1j)", f"{var99:.2f}%"),
        ("CVaR 95%", f"{cvar95:.2f}%"), ("Downside Risk", f"{downside:.2f}%"),
    ]
    for i, (lbl, val) in enumerate(risk_data):
        _add_card(s, 0.2 + i * 2.4, 0.7, 2.2, 0.9, lbl, val, DOWN)
    _add_rect(s, 0.2, 1.8, 9.6, 1.3, CARD)
    _add_rect(s, 0.2, 1.8, 0.06, 1.3, ORA)
    _add_text(s,
        f"Avec un niveau de confiance de 95%, la perte maximale attendue sur une journée est de "
        f"{abs(var95):.2f}%. Dans les 5% de pires cas (CVaR 95%), la perte moyenne serait de {abs(cvar95):.2f}%.",
        0.4, 1.88, 9.2, 1.1, font_size=12, color=WHITE)
    _footer(s)

    # ── SLIDE 10 : Disclaimer ──
    s = prs.slides.add_slide(blank)
    _set_bg(s, CARD)
    _add_rect(s, 0, 0, 0.12, 5.625, ACC)
    _add_text(s, "Avertissement", 0.4, 1.0, 9.2, 0.7,
              font_size=22, bold=True, color=ORA)
    _add_text(s,
        "Ce document est généré à titre éducatif uniquement par Apex Markets.\n"
        "Les données proviennent de yfinance et peuvent être décalées ou incomplètes.\n"
        "Aucune information contenue dans ce support ne constitue une recommandation d'investissement.\n"
        "Tout investissement comporte un risque de perte en capital.",
        0.4, 1.9, 9.2, 2.2, font_size=13, color=MUTED)
    _add_text(s, f"Apex Markets — {datetime.now().strftime('%d/%m/%Y')}",
              0.4, 4.8, 9.2, 0.35, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Sérialiser ──
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    main()
